#!/usr/bin/env python3
"""Merge slides from one PowerPoint deck into another.

`rearrange.py` only reorders/duplicates slides WITHIN a single file. This script
fills the gap: it copies slides from a source deck (e.g. a reusable module) into
a base deck, importing each slide's shapes, slide-level background, and embedded
images, and attaching them to a layout in the base deck so they inherit the base
deck's master (footer, theme).

It is built for the common case where base deck and modules share a theme, so
attaching imported slides to a base layout re-brands them consistently (e.g. the
base deck's footer) while their own explicit shapes carry the content.

Usage:
    # Append every slide of a module to the end of the base deck.
    python merge_decks.py base.pptx output.pptx --add module.pptx

    # Insert specific source slides (1-based) at position 12 (0-based) in base.
    python merge_decks.py base.pptx output.pptx --add module.pptx --slides 1,2,3 --at 12

    # Pin imported slides to a specific base layout (see --list-layouts).
    python merge_decks.py base.pptx output.pptx --add module.pptx --layout 5

    # Inspect base layouts to choose one.
    python merge_decks.py base.pptx --list-layouts

Notes:
    - Source slide numbers are 1-based; --at is a 0-based index into the base
      deck's slide list (0 = before the first slide).
    - Images and external hyperlinks are carried over. Charts, embedded OLE
      objects, and other rare relationship types are NOT copied — the script
      warns if it sees one so you can handle that slide manually.
    - Imported slides inherit the base deck's master/footer by design. If a
      source slide set its own slide-level background, that background is kept.
"""

import argparse
import copy
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

IMAGE_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
HYPERLINK_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"
# Relationship roles that the base layout/master already provides; ignore them.
STRUCTURAL_RELTYPES = {
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
}
# r: attributes that may reference a relationship inside shape XML.
EMBED_ATTRS = [qn("r:embed"), qn("r:link"), qn("r:id"), qn("r:pict")]


def list_layouts(prs):
    print(f"{len(prs.slide_layouts)} layouts on the base deck's primary master:")
    for i, layout in enumerate(prs.slide_layouts):
        shape_count = len(layout.shapes)
        print(f"  [{i:>2}] {layout.name}  ({shape_count} shapes)")


def pick_default_layout(prs):
    """Prefer a 'blank'-style layout; otherwise the one with the fewest shapes."""
    layouts = list(prs.slide_layouts)
    by_name = [l for l in layouts if "blank" in l.name.lower()]
    if by_name:
        return by_name[0]
    return min(layouts, key=lambda l: len(l.shapes))


def strip_layout_shapes(slide):
    """Remove the placeholder shapes the layout seeded onto a fresh slide, so only
    imported content (plus inherited layout/master visuals) remains."""
    spTree = slide.shapes._spTree
    for shape in list(slide.shapes):
        spTree.remove(shape.element)


def copy_background(src_slide, dest_slide):
    """Carry over a slide-level <p:bg> if the source set one."""
    src_cSld = src_slide._element.find(qn("p:cSld"))
    dest_cSld = dest_slide._element.find(qn("p:cSld"))
    if src_cSld is None or dest_cSld is None:
        return
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is None:
        return
    # Remove any existing bg on the destination, then insert as first child of cSld.
    existing = dest_cSld.find(qn("p:bg"))
    if existing is not None:
        dest_cSld.remove(existing)
    dest_cSld.insert(0, copy.deepcopy(src_bg))


def import_relationships(src_slide, dest_slide, warnings):
    """Recreate the source slide's image and hyperlink relationships on the
    destination slide. Returns {old_rId: new_rId} for rewriting shape XML."""
    rid_map = {}
    tmp_files = []
    for rId, rel in src_slide.part.rels.items():
        if rel.reltype in STRUCTURAL_RELTYPES:
            continue
        if rel.reltype == IMAGE_RELTYPE:
            image_part = rel.target_part
            ext = image_part.partname.ext  # e.g. 'png'
            tmp = tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False)
            tmp.write(image_part.blob)
            tmp.close()
            tmp_files.append(tmp.name)
            _, new_rid = dest_slide.part.get_or_add_image_part(tmp.name)
            rid_map[rId] = new_rid
        elif rel.reltype == HYPERLINK_RELTYPE and rel.is_external:
            new_rid = dest_slide.part.relate_to(
                rel.target_ref, rel.reltype, is_external=True
            )
            rid_map[rId] = new_rid
        else:
            warnings.append(
                f"unsupported relationship {rel.reltype.split('/')[-1]} "
                f"(rId {rId}) was not copied — review this slide manually"
            )
    return rid_map, tmp_files


def rewrite_rel_ids(element, rid_map):
    """Point every r:embed / r:link / r:id in the copied XML at its new rId."""
    for el in element.iter():
        for attr in EMBED_ATTRS:
            old = el.get(attr)
            if old is not None and old in rid_map:
                el.set(attr, rid_map[old])


def copy_slide(src_slide, dest_prs, dest_layout, warnings):
    """Create a new slide in dest_prs from src_slide's content."""
    new_slide = dest_prs.slides.add_slide(dest_layout)
    strip_layout_shapes(new_slide)
    copy_background(src_slide, new_slide)

    rid_map, tmp_files = import_relationships(src_slide, new_slide, warnings)

    # Deep-copy every shape from the source spTree (skip the group's own
    # nvGrpSpPr / grpSpPr bookkeeping elements).
    src_spTree = src_slide.shapes._spTree
    dest_spTree = new_slide.shapes._spTree
    skip = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}
    for child in list(src_spTree):
        if child.tag in skip:
            continue
        new_el = copy.deepcopy(child)
        rewrite_rel_ids(new_el, rid_map)
        dest_spTree.append(new_el)

    for path in tmp_files:
        try:
            Path(path).unlink()
        except OSError:
            pass
    return new_slide


def move_slide_to(dest_prs, from_index, to_index):
    """Move the slide currently at from_index to to_index in the slide list."""
    sldIdLst = dest_prs.slides._sldIdLst
    elements = list(sldIdLst)
    el = elements[from_index]
    sldIdLst.remove(el)
    sldIdLst.insert(to_index, el)


def merge(base_path, output_path, module_path, slide_nums, at_index, layout_index):
    dest_prs = Presentation(str(base_path))
    src_prs = Presentation(str(module_path))

    total_src = len(src_prs.slides)
    if slide_nums:
        for n in slide_nums:
            if n < 1 or n > total_src:
                sys.exit(f"Error: source slide {n} out of range (1-{total_src})")
        chosen = slide_nums
    else:
        chosen = list(range(1, total_src + 1))

    if layout_index is not None:
        if layout_index < 0 or layout_index >= len(dest_prs.slide_layouts):
            sys.exit(f"Error: layout {layout_index} out of range "
                     f"(0-{len(dest_prs.slide_layouts) - 1}). Use --list-layouts.")
        dest_layout = dest_prs.slide_layouts[layout_index]
    else:
        dest_layout = pick_default_layout(dest_prs)

    base_count = len(dest_prs.slides)
    if at_index is None:
        at_index = base_count
    if at_index < 0 or at_index > base_count:
        sys.exit(f"Error: --at {at_index} out of range (0-{base_count})")

    all_warnings = []
    print(f"Merging {len(chosen)} slide(s) from {module_path.name} into {base_path.name}")
    print(f"  attaching to base layout: '{dest_layout.name}'")

    # add_slide always appends; we copy in order, then move the appended block
    # into position so source order is preserved at the insertion point.
    for offset, n in enumerate(chosen):
        warnings = []
        copy_slide(src_prs.slides[n - 1], dest_prs, dest_layout, warnings)
        appended_index = len(dest_prs.slides) - 1
        move_slide_to(dest_prs, appended_index, at_index + offset)
        tag = f"  [src {n}] -> base position {at_index + offset}"
        if warnings:
            tag += "  ⚠ " + "; ".join(warnings)
            all_warnings.extend(warnings)
        print(tag)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    dest_prs.save(str(output_path))
    print(f"\nSaved: {output_path}  ({len(dest_prs.slides)} slides total)")
    if all_warnings:
        print(f"\n{len(all_warnings)} warning(s) — review the flagged slides in PowerPoint.")


def main():
    parser = argparse.ArgumentParser(
        description="Merge slides from a module deck into a base deck.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("base", help="Base .pptx (the working deck)")
    parser.add_argument("output", nargs="?", help="Output .pptx (omit for --list-layouts)")
    parser.add_argument("--add", help="Source module .pptx to merge in")
    parser.add_argument("--slides", help="Comma-separated source slide numbers (1-based); default all")
    parser.add_argument("--at", type=int, help="0-based insert position in base; default append")
    parser.add_argument("--layout", type=int, help="Base layout index to attach imported slides to")
    parser.add_argument("--list-layouts", action="store_true", help="List base layouts and exit")
    args = parser.parse_args()

    base_path = Path(args.base)
    if not base_path.exists():
        sys.exit(f"Error: base file not found: {base_path}")

    if args.list_layouts:
        list_layouts(Presentation(str(base_path)))
        return

    if not args.add:
        sys.exit("Error: --add <module.pptx> is required (or use --list-layouts)")
    module_path = Path(args.add)
    if not module_path.exists():
        sys.exit(f"Error: module file not found: {module_path}")
    if not args.output:
        sys.exit("Error: output .pptx path required")

    slide_nums = None
    if args.slides:
        try:
            slide_nums = [int(x.strip()) for x in args.slides.split(",")]
        except ValueError:
            sys.exit("Error: --slides must be comma-separated integers (e.g. 1,2,3)")

    merge(base_path, Path(args.output), module_path, slide_nums, args.at, args.layout)


if __name__ == "__main__":
    main()
