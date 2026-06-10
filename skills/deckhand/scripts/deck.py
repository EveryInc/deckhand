#!/usr/bin/env python3
"""deck.py — unified deck editor. The LLM writes patches; this executes them.

One addressing scheme for EVERY shape (text, picture, table, group): the shape's
native id, written `s<id>` (stable — deleting a shape never renumbers others).
All slide indices are 0-BASED everywhere in this tool (slide 0 = first slide).

Subcommands
    docs      print the full patch reference + semantics (`deck.py docs`, no file needed)
    inspect   read everything as JSON: geometry, text+fonts, images (rId+media),
              fill/line colors, issues; --master adds masters/layouts (footers)
    apply     execute a JSON patch of ops atomically; lint after save
    fix       deterministic geometry repair (overflow/off-slide); reports residue
    render    slides → JPGs named slide-<index>.jpg (0-based, hidden-slide aware);
              --crop l,t,w,h --scale 2 zooms a region (inches, matches inspect)
    diff      structural changelog between two decks (no rendering needed)
    slides    reorder/duplicate/delete slides by 0-based sequence (0,3,3,5)
    merge     copy slides from another deck into this one (--slides --at --layout)
    xml       escape hatch: get pretty-printed part XML / set it back validated

Usage
    python3 deck.py deck.pptx inspect [--slide 3,7] [--issues] [-o out.json]
    python3 deck.py deck.pptx apply patch.json -o out.pptx
    python3 deck.py deck.pptx fix --slides 3,7 -o out.pptx
    python3 deck.py deck.pptx render -o imgdir/ [--slide 3,7] [--dpi 110]
    python3 deck.py deck.pptx diff other.pptx
    python3 deck.py deck.pptx xml get --slide 5 -o slide5.xml
    python3 deck.py deck.pptx xml set slide5.xml --slide 5 -o out.pptx

Patch format (apply): {"ops": [ ... ]} — ops run in order, all-or-nothing.
    {"op":"set-text",    "slide":3, "shape":"s12", "text":["Title", "Subtitle"]}
        Each new paragraph INHERITS the formatting of the paragraph it replaces
        (extras inherit from the last). Pass an object instead of a string to
        override: {"text":"Big", "font_size":28, "bold":true}. Table cells:
        add "cell":[row,col].
    {"op":"swap-image",  "slide":3, "rid":"rId3", "image":"/abs/new.png"}
        or "shape":"s9". "fit":"auto" (default) keeps aspect ratio when it
        clearly differs from the frame; "stretch"/"contain" force behaviour.
        Global: {"op":"swap-image","media":"image13.png","image":"/abs/new.png"}
        overwrites the media bytes — every slide/master using it changes.
    {"op":"replace-text","scope":"deck",   "from":"Globex Corp", "to":"Acme"}
    {"op":"replace-text","scope":"master", "from":"Globex Cowork", "to":"Acme Cowork"}
        scope: "deck" (all slides) | "master" (masters+layouts) | "slide" + "slide":N
    {"op":"set-notes",   "slide":3, "notes":"speaker notes text"}
    {"op":"move",        "slide":3, "shape":"s12", "to":[1.0,2.5]}   # inches; or "by":[dx,dy]
    {"op":"resize",      "slide":3, "shape":"s12", "size":[4.0,1.5]} # inches; or "scale":0.8
    {"op":"set-style",   "slide":3, "shape":"s12", "font_size":18, "color":"FFFFFF",
                         "bold":true, "font_name":"Arial", "fill":"0B3D3A"}
    {"op":"delete",      "slide":3, "shape":"s12"}
    {"op":"duplicate",   "slide":3, "shape":"s12", "offset":[0,1.2], "text":["Fourth pillar"]}
    {"op":"copy-shape",  "from_slide":8, "shape":"s12", "slide":3, "at":[1.0,2.0], "text":["…"]}
    {"op":"add-shape",   "slide":3, "kind":"textbox", "at":[1,2], "size":[4,1.5], "text":["…"]}
        kinds: textbox, rect, rounded_rect, ellipse, line (from/to), any MSO_SHAPE
        name; takes every set-style key (fill, line_color, font_size, rotation…)
    {"op":"add-picture", "slide":3, "image":"/abs/img.png", "at":[1,2], "width":4}
    {"op":"add-table",   "slide":3, "at":[1,2], "size":[8,3], "rows":[["A","B"],["1","2"]]}
    {"op":"add-slide",   "layout":"Blank", "at":5}   # layout name or index; omit = blank-est
    {"op":"reorder",     "slide":3, "shape":"s12", "z":"front"}  # back|forward|backward
    {"op":"add-row",     "slide":3, "shape":"s12", "cells":["a","b"], "copy":-1, "at":2}
    {"op":"delete-row",  "slide":3, "shape":"s12", "row":2}    # add-col/delete-col likewise
"""

from __future__ import annotations

import argparse
import contextlib
import copy
import io
import json
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

LINE_DASH_STYLES = {
    "solid": MSO_LINE_DASH_STYLE.SOLID,
    "dash": MSO_LINE_DASH_STYLE.DASH,
    "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
    "dash_dot_dot": MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
    "long_dash": MSO_LINE_DASH_STYLE.LONG_DASH,
    "long_dash_dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
    "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
    "round_dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
    "square_dot": MSO_LINE_DASH_STYLE.SQUARE_DOT,
}

from inventory import ParagraphData, ShapeData
from replace import apply_font_properties, apply_paragraph_properties

EMU_PER_IN = 914400
IMAGE_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
R_EMBED = qn("r:embed")
R_LINK = qn("r:link")
R_ID = qn("r:id")
GROUP_TYPE = 6  # MSO_SHAPE_TYPE.GROUP


def inches(emu):
    return round(emu / EMU_PER_IN, 2)


# ---------------------------------------------------------------------------
# Shape index — the one addressing scheme
# ---------------------------------------------------------------------------

class TF:
    """Maps child-space EMU coords of a container to absolute EMU.
    abs_x = ox + (x - chx) * sx"""

    def __init__(self, ox=0, oy=0, sx=1.0, sy=1.0, chx=0, chy=0):
        self.ox, self.oy, self.sx, self.sy, self.chx, self.chy = ox, oy, sx, sy, chx, chy

    def abs_x(self, x):
        return self.ox + (x - self.chx) * self.sx

    def abs_y(self, y):
        return self.oy + (y - self.chy) * self.sy

    def child_x(self, ax):
        return self.chx + (ax - self.ox) / self.sx

    def child_y(self, ay):
        return self.chy + (ay - self.oy) / self.sy


class Rec:
    """One shape's record in the index."""

    def __init__(self, slide_idx, shape, tf, group_sid=None):
        self.slide_idx = slide_idx
        self.shape = shape
        self.tf = tf
        self.group = group_sid
        self.sid = "s%d" % shape.shape_id
        self.name = shape.name
        try:
            st = shape.shape_type
            self.type = "UNKNOWN" if st is None else str(st).split(".")[-1].split(" ")[0]
        except Exception:
            self.type = "UNKNOWN"
        l = shape.left or 0
        t = shape.top or 0
        self.left = inches(tf.abs_x(l))
        self.top = inches(tf.abs_y(t))
        self.width = inches((shape.width or 0) * tf.sx)
        self.height = inches((shape.height or 0) * tf.sy)
        self.is_text = bool(
            getattr(shape, "has_text_frame", False)
            and shape.has_text_frame
            and shape.text_frame.text.strip()
        )
        self.is_table = bool(getattr(shape, "has_table", False) and shape.has_table)
        # picture blips: [(rid, media_filename)]
        self.rids = []
        if self.type == "PICTURE" or (not self.is_table and self.type != "GROUP"):
            part = shape.part
            for blip in shape._element.findall(".//" + qn("a:blip")):
                rid = blip.get(R_EMBED)
                if rid and rid in part.rels:
                    self.rids.append((rid, part.rels[rid].target_ref.split("/")[-1]))
        self.sd = None  # ShapeData, attached when measuring
        self.covered_by = {}  # {picture_sid: sq in} — pictures above this text in z-order

    def text_preview(self, n=48):
        if self.is_text:
            t = " / ".join(
                p.text.strip() for p in self.shape.text_frame.paragraphs if p.text.strip()
            )
            return (t[: n - 1] + "…") if len(t) > n else t
        if self.is_table:
            return "(table)"
        if self.rids:
            return "(image %s)" % self.rids[0][1]
        return ""


def _group_child_tf(gshape, tf):
    """TF for the children of a group, honoring chOff/chExt scaling."""
    g_abs_l = tf.abs_x(gshape.left or 0)
    g_abs_t = tf.abs_y(gshape.top or 0)
    g_abs_w = (gshape.width or 0) * tf.sx
    g_abs_h = (gshape.height or 0) * tf.sy
    chx, chy = gshape.left or 0, gshape.top or 0
    chw, chh = gshape.width or 0, gshape.height or 0
    xfrm = gshape._element.find(qn("p:grpSpPr") + "/" + qn("a:xfrm"))
    if xfrm is not None:
        choff = xfrm.find(qn("a:chOff"))
        chext = xfrm.find(qn("a:chExt"))
        if choff is not None:
            chx, chy = int(choff.get("x")), int(choff.get("y"))
        if chext is not None:
            chw, chh = int(chext.get("cx")), int(chext.get("cy"))
    sx = (g_abs_w / chw) if chw else 1.0
    sy = (g_abs_h / chh) if chh else 1.0
    return TF(g_abs_l, g_abs_t, sx, sy, chx, chy)


def build_index(prs, measure=False, only_slides=None):
    """{slide_idx: {sid: Rec}}. measure=True attaches ShapeData (NOTE: measuring
    touches font properties and may add empty elements — never measure the
    instance you intend to save)."""
    index = {}
    for slide_idx, slide in enumerate(prs.slides):
        if only_slides is not None and slide_idx not in only_slides:
            continue
        recs = {}

        def walk(shapes, tf, group_sid=None):
            for sh in shapes:
                try:
                    is_group = sh.shape_type == GROUP_TYPE
                except Exception:
                    is_group = False
                rec = Rec(slide_idx, sh, tf, group_sid)
                if rec.sid in recs:  # duplicate native id (spec violation) — disambiguate
                    n = 2
                    while "%s-%d" % (rec.sid, n) in recs:
                        n += 1
                    rec.sid = "%s-%d" % (rec.sid, n)
                recs[rec.sid] = rec
                if is_group:
                    walk(sh.shapes, _group_child_tf(sh, tf), rec.sid)

        walk(slide.shapes, TF())
        if measure:
            # measure ALL shapes (slide-overflow applies to pictures too);
            # frame-overflow only fires on shapes with text
            for r in recs.values():
                r.sd = ShapeData(
                    r.shape,
                    int(r.left * EMU_PER_IN),
                    int(r.top * EMU_PER_IN),
                    slide,
                )
                r.sd.shape_id = r.sid
            # overlap detection between text shapes only (text-over-background is normal)
            sds = [r.sd for r in recs.values() if r.is_text]
            if len(sds) > 1:
                from inventory import detect_overlaps

                detect_overlaps(sds)
            # text drawn BELOW a picture in z-order renders clipped/hidden
            # behind it — include estimated overflow, since a re-wrapped last
            # line that slides under a picture is invisible in thumbnails.
            # (Text drawn ON TOP of a picture is normal design; not flagged.)
            from inventory import calculate_overlap

            rec_list = list(recs.values())
            for i, r in enumerate(rec_list):
                if not r.is_text or r.sd is None:
                    continue
                ov = r.sd.frame_overflow_bottom or 0
                trect = (r.left, r.top, r.width, r.height + ov)
                for o in rec_list[i + 1:]:
                    if o.type != "PICTURE":
                        continue
                    hit, area = calculate_overlap(trect, (o.left, o.top, o.width, o.height))
                    if hit:
                        r.covered_by[o.sid] = area
        index[slide_idx] = recs
    return index


def rec_issues(rec):
    """Geometric issues for one measured Rec, or {}.

    Non-text shapes that bleed off MULTIPLE edges or cover most of the slide
    are treated as intentional design (statues, textures) and not flagged."""
    if rec.sd is None:
        return {}
    out = {}
    if rec.sd.frame_overflow_bottom is not None:
        out["frame_overflow_bottom"] = rec.sd.frame_overflow_bottom
    off_edges = sum(
        1 for v in (rec.sd.slide_overflow_right, rec.sd.slide_overflow_bottom) if v is not None
    )
    bleed = False
    if not rec.is_text and rec.sd.slide_width_emu:
        slide_area = inches(rec.sd.slide_width_emu) * inches(rec.sd.slide_height_emu)
        bleed = off_edges >= 2 or (rec.width * rec.height) > 0.6 * slide_area
    if not bleed:
        if rec.sd.slide_overflow_right is not None:
            out["slide_overflow_right"] = rec.sd.slide_overflow_right
        if rec.sd.slide_overflow_bottom is not None:
            out["slide_overflow_bottom"] = rec.sd.slide_overflow_bottom
    if rec.sd.overlapping_shapes:
        out["overlaps"] = rec.sd.overlapping_shapes
    if rec.covered_by:
        out["covered_by"] = rec.covered_by
    if rec.sd.warnings:
        out["warnings"] = rec.sd.warnings
    return out


def slide_listing(index, slide_idx):
    """Human-readable shape listing — used in error messages so the agent can
    self-correct without re-running inspect."""
    recs = index.get(slide_idx, {})
    lines = ["shapes on slide %d:" % slide_idx]
    for sid, r in recs.items():
        extra = " in %s" % r.group if r.group else ""
        lines.append(
            "  %-6s %-12s [%.2f,%.2f %sx%sin]%s  %s"
            % (sid, r.type, r.left, r.top, r.width, r.height, extra, r.text_preview())
        )
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# inspect
# ---------------------------------------------------------------------------

def _run_dict(run):
    """One run's text + non-default font props (same vocabulary set-text accepts)."""
    d = {"text": run.text}
    f = run.font
    if f.name:
        d["font_name"] = f.name
    if f.size:
        d["font_size"] = f.size.pt
    if f.bold is not None:
        d["bold"] = f.bold
    if f.italic is not None:
        d["italic"] = f.italic
    if f.underline is not None:
        d["underline"] = f.underline
    try:
        if f.color.rgb:
            d["color"] = str(f.color.rgb)
    except (AttributeError, TypeError):
        try:
            if f.color.theme_color:
                d["theme_color"] = f.color.theme_color.name
        except (AttributeError, TypeError, ValueError):
            pass
    try:
        if run.hyperlink.address:
            d["link"] = run.hyperlink.address
    except (AttributeError, KeyError):
        pass
    return d


def _para_dict_with_runs(p):
    """ParagraphData dict, plus a per-run breakdown when the paragraph mixes
    run formats (so agents can SEE a mid-sentence bold word and patch it with
    set-text "runs")."""
    d = ParagraphData(p).to_dict()
    runs = [r for r in p.runs if r.text]
    if len(runs) > 1:
        rdicts = [_run_dict(r) for r in runs]
        fmts = [{k: v for k, v in rd.items() if k != "text"} for rd in rdicts]
        if any(f != fmts[0] for f in fmts[1:]):
            d["runs"] = rdicts
    elif runs:
        link = _run_dict(runs[0]).get("link")
        if link:
            d["link"] = link
    return d


def _shape_colors(shape):
    """Solid fill + line color/width when explicitly set. Only call on
    instances that will never be saved (color access can normalize XML)."""
    out = {}
    try:
        f = shape.fill
        if f.type is not None and str(f.type).startswith("SOLID"):
            try:
                out["fill"] = str(f.fore_color.rgb)
            except Exception:
                try:
                    out["fill"] = "theme:" + f.fore_color.theme_color.name
                except Exception:
                    pass
        elif f.type is not None and str(f.type).startswith("GRADIENT"):
            g = {}
            try:
                g["stops"] = [
                    {"color": str(s.color.rgb), "position": round(s.position, 3)}
                    for s in f.gradient_stops
                ]
            except Exception:
                pass
            try:
                g["angle"] = round(f.gradient_angle, 1)
            except Exception:
                pass
            if g:
                out["fill_gradient"] = g
    except Exception:
        pass
    try:
        ln = shape.line
        d = {}
        try:
            if ln.width:
                d["width_pt"] = round(ln.width.pt, 2)
        except Exception:
            pass
        try:
            d["color"] = str(ln.color.rgb)
        except Exception:
            try:
                d["color"] = "theme:" + ln.color.theme_color.name
            except Exception:
                pass
        if d:
            out["line"] = d
    except Exception:
        pass
    return out


def _slide_fonts(slide_like):
    return sorted(
        {el.get("typeface") for el in slide_like._element.iter(qn("a:latin")) if el.get("typeface")}
    )


def _master_entry(part_obj):
    """Light-weight shape listing for a master or layout (footer edits target these)."""
    recs = {}

    def walk(shapes, tf, group_sid=None):
        for sh in shapes:
            try:
                is_group = sh.shape_type == GROUP_TYPE
            except Exception:
                is_group = False
            r = Rec(0, sh, tf, group_sid)
            recs[r.sid] = r
            if is_group:
                walk(sh.shapes, _group_child_tf(sh, tf), r.sid)

    walk(part_obj.shapes, TF())
    out = {}
    for sid, r in recs.items():
        e = {"name": r.name, "type": r.type, "pos": [r.left, r.top], "size": [r.width, r.height]}
        if r.is_text:
            e["text"] = [p.text for p in r.shape.text_frame.paragraphs if p.text.strip()]
        if r.rids:
            e["media"] = r.rids[0][1]
        out[sid] = e
    fonts = _slide_fonts(part_obj)
    if fonts:
        out["_fonts"] = fonts
    return out


def cmd_inspect(args):
    prs = Presentation(args.file)
    only = parse_slide_list(args.slide) if args.slide else None
    index = build_index(prs, measure=True, only_slides=only)

    if getattr(args, "brief", False):
        # compact orientation view: one line per shape; full JSON only when
        # you're about to write a patch
        lines = []
        for slide_idx, recs in index.items():
            shown = []
            for sid, r in recs.items():
                issues = rec_issues(r)
                if args.issues and not issues:
                    continue
                bits = ["%-6s %-11s [%.2f,%.2f %sx%s]" % (sid, r.type, r.left, r.top, r.width, r.height)]
                if r.group:
                    bits.append("in " + r.group)
                if r.is_text:
                    p0 = next((p for p in r.shape.text_frame.paragraphs if p.text.strip()), None)
                    if p0 is not None:
                        pd = ParagraphData(p0)
                        f = []
                        if pd.font_size:
                            f.append("%gpt" % pd.font_size)
                        if pd.font_name:
                            f.append(pd.font_name)
                        if f:
                            bits.append(" ".join(f))
                pv = r.text_preview(44)
                if pv:
                    bits.append(pv)
                if issues:
                    bits.append("⚠ " + ", ".join("%s=%s" % (k, v) for k, v in issues.items()))
                shown.append("  " + "  ".join(bits))
            if shown:
                fonts = _slide_fonts(prs.slides[slide_idx])
                lines.append("slide %d%s:" % (slide_idx, (" (fonts: %s)" % ", ".join(fonts)) if fonts else ""))
                lines.extend(shown)
        text = "\n".join(lines) if lines else ("(no shapes with issues)" if args.issues else "(no shapes)")
        if args.output:
            Path(args.output).write_text(text + "\n", encoding="utf-8")
            print("Wrote %s" % args.output)
        else:
            print(text)
        return
    out = {
        "file": str(args.file),
        "slide_count": len(prs.slides._sldIdLst),
        "slide_size": [
            inches(prs.slide_width),
            inches(prs.slide_height),
        ],
        "slides": {},
    }
    for slide_idx, recs in index.items():
        slide_out = {}
        for sid, r in recs.items():
            issues = rec_issues(r)
            if args.issues and not issues:
                continue
            entry = {
                "name": r.name,
                "type": r.type,
                "pos": [r.left, r.top],
                "size": [r.width, r.height],
            }
            if r.group:
                entry["group"] = r.group
            if r.rids:
                entry["rid"] = r.rids[0][0]
                entry["media"] = r.rids[0][1]
                if len(r.rids) > 1:
                    entry["all_rids"] = [{"rid": a, "media": b} for a, b in r.rids]
            if r.is_text:
                if r.sd is not None and r.sd.placeholder_type:
                    entry["placeholder"] = r.sd.placeholder_type
                entry["paragraphs"] = [
                    _para_dict_with_runs(p)
                    for p in r.shape.text_frame.paragraphs
                    if p.text.strip()
                ]
            if r.is_table:
                entry["rows"] = [
                    [c.text for c in row.cells] for row in r.shape.table.rows
                ]
            if r.type != "GROUP":
                entry.update(_shape_colors(r.shape))
            alt = _get_alt_text(r.shape)
            if alt:
                entry["alt_text"] = alt
            try:
                if r.shape.rotation:
                    entry["rotation"] = round(r.shape.rotation, 1)
            except (AttributeError, ValueError):
                pass
            if issues:
                entry["issues"] = issues
            slide_out[sid] = entry
        if slide_out and not args.issues:
            fonts = _slide_fonts(prs.slides[slide_idx])
            if fonts:
                slide_out["_fonts"] = fonts
            s = prs.slides[slide_idx]
            if s.has_notes_slide:
                notes = s.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_out["_notes"] = notes
            tr = _transition_dict(s)
            if tr:
                slide_out["_transition"] = tr
        if slide_out or not args.issues:
            if not args.issues or slide_out:
                out["slides"][str(slide_idx)] = slide_out
    hidden = [i for i, s in enumerate(prs.slides) if s.element.get("show") == "0"]
    if hidden:
        out["hidden_slides"] = hidden
    cp = prs.core_properties
    props = {k: getattr(cp, k) for k in PROPS_KEYS if getattr(cp, k, None)}
    if props:
        out["props"] = props
    if getattr(args, "master", False):
        masters = {}
        for mi, master in enumerate(prs.slide_masters):
            masters["master-%d" % mi] = _master_entry(master)
            for li, layout in enumerate(master.slide_layouts):
                masters["master-%d/layout-%d (%s)" % (mi, li, layout.name)] = _master_entry(layout)
        out["masters"] = masters
    text = json.dumps(out, indent=2, ensure_ascii=False)
    if args.output:
        Path(args.output).write_text(text, encoding="utf-8")
        print("Wrote %s (%d slides)" % (args.output, len(out["slides"])))
    else:
        print(text)


# ---------------------------------------------------------------------------
# apply — the patch engine
# ---------------------------------------------------------------------------

VALID_OPS = (
    "set-text swap-image replace-text replace-color set-notes set-props "
    "set-slide set-theme move resize set-style "
    "delete duplicate copy-shape "
    "add-shape add-picture add-table add-slide reorder "
    "add-row delete-row add-col delete-col"
).split()

# friendly aliases; any other MSO_SHAPE name (e.g. "CHEVRON") also works
SHAPE_KINDS = {
    "rect": MSO_SHAPE.RECTANGLE,
    "rounded_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "ellipse": MSO_SHAPE.OVAL,
    "oval": MSO_SHAPE.OVAL,
}


def resolve_shape_kind(kind):
    """'textbox' | 'line' | an MSO_SHAPE for everything else."""
    if kind in ("textbox", "line"):
        return kind
    if kind in SHAPE_KINDS:
        return SHAPE_KINDS[kind]
    try:
        return MSO_SHAPE[kind.upper().replace("-", "_")]
    except KeyError:
        raise PatchError(
            "unknown shape kind '%s' — use textbox, line, %s, or any MSO_SHAPE name (e.g. CHEVRON, PENTAGON)"
            % (kind, ", ".join(sorted(SHAPE_KINDS)))
        )


class Ctx:
    def __init__(self, prs, index):
        self.prs = prs
        self.index = index
        self.touched = set()  # slide indices modified
        self.log = []

    def rec(self, slide_idx, shape_ref):
        recs = self.index.get(slide_idx)
        if recs is None:
            raise PatchError("slide %s out of range (0-%d)" % (slide_idx, len(self.prs.slides._sldIdLst) - 1))
        if shape_ref in recs:
            return recs[shape_ref]
        # fall back to exact shape-name match
        named = [r for r in recs.values() if r.name == shape_ref]
        if len(named) == 1:
            return named[0]
        raise PatchError(
            "shape '%s' not found on slide %d.\n%s" % (shape_ref, slide_idx, slide_listing(self.index, slide_idx))
        )

    def reindex_slide(self, slide_idx):
        sub = build_index(self.prs, measure=False, only_slides={slide_idx})
        self.index[slide_idx] = sub[slide_idx]


class PatchError(Exception):
    pass


def parse_slide_list(s):
    return {int(x.strip()) for x in str(s).split(",") if x.strip() != ""}


# font keys that may be overridden per RUN (everything else is paragraph-level)
RUN_FONT_KEYS = ("bold", "italic", "underline", "font_size", "font_name", "color", "theme_color", "link")


def _merge_font(base, override):
    """Merge run font overrides onto an inherited base. An explicit color kind
    (rgb vs theme) in the override displaces the other kind from the base."""
    merged = dict(base)
    if "color" in override:
        merged.pop("theme_color", None)
    if "theme_color" in override:
        merged.pop("color", None)
    merged.update(override)
    return merged


def norm_paragraph_items(value):
    items = value if isinstance(value, list) else [value]
    out = []
    for it in items:
        if isinstance(it, str):
            out.append({"text": it})
        elif isinstance(it, dict) and "runs" in it:
            if "text" in it:
                raise PatchError("a paragraph takes 'text' OR 'runs', not both")
            runs = it["runs"]
            if not isinstance(runs, list) or not runs:
                raise PatchError('"runs" must be a non-empty array of {"text": ..., <font overrides>}')
            norm_runs = []
            for r in runs:
                if isinstance(r, str):
                    norm_runs.append({"text": r})
                elif isinstance(r, dict) and "text" in r:
                    bad = sorted(k for k in r if k != "text" and k not in RUN_FONT_KEYS)
                    if bad:
                        raise PatchError(
                            "run key(s) %s not supported — per-run keys are: %s. "
                            "Paragraph-level keys (alignment, bullet, level, spacing) go on the paragraph object."
                            % (", ".join(bad), ", ".join(RUN_FONT_KEYS))
                        )
                    norm_runs.append(dict(r))
                else:
                    raise PatchError("each run must be a string or an object with 'text'")
            d = dict(it)
            d["runs"] = norm_runs
            out.append(d)
        elif isinstance(it, dict) and "text" in it:
            out.append(dict(it))
        else:
            raise PatchError("each text item must be a string or an object with 'text' (or 'runs')")
    return out


def write_paragraphs_inherit(tf, items):
    """Clear a text frame and write new paragraphs, each inheriting the
    formatting of the paragraph it replaces (extras inherit from the last).
    A paragraph item may carry "runs" for mixed in-paragraph formatting: every
    run inherits the OLD paragraph's first-run font (the same baseline plain
    set-text uses), then the run's own keys override."""
    old = [ParagraphData(p).to_dict() for p in tf.paragraphs if p.text.strip()]
    tf.clear()
    for i, item in enumerate(items):
        base = dict(old[min(i, len(old) - 1)]) if old else {}
        base.pop("text", None)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        runs = item.get("runs")
        if not runs:
            merged = dict(base)
            merged.update(item)
            apply_paragraph_properties(p, merged)
            continue
        para = dict(base)
        para.update({k: v for k, v in item.items() if k != "runs"})
        font_base = {k: para[k] for k in RUN_FONT_KEYS if k in para}
        # paragraph props + first run in one go
        first = {k: v for k, v in para.items() if k not in RUN_FONT_KEYS}
        first.update(_merge_font(font_base, runs[0]))
        first["text"] = runs[0]["text"]
        apply_paragraph_properties(p, first)
        for rspec in runs[1:]:
            run = p.add_run()
            run.text = rspec["text"]
            apply_font_properties(run, _merge_font(font_base, rspec))


def op_set_text(ctx, op):
    rec = ctx.rec(op["slide"], op["shape"])
    items = norm_paragraph_items(op["text"])
    if "cell" in op:
        if not rec.is_table:
            raise PatchError("'cell' given but %s on slide %d is not a table" % (rec.sid, op["slide"]))
        r, c = op["cell"]
        try:
            tf = rec.shape.table.cell(r, c).text_frame
        except IndexError:
            t = rec.shape.table
            raise PatchError(
                "cell [%d,%d] out of range — table is %d rows x %d cols"
                % (r, c, len(t.rows), len(t.columns))
            )
    elif rec.is_table:
        raise PatchError("%s is a table — address cells with \"cell\":[row,col]" % rec.sid)
    elif not getattr(rec.shape, "has_text_frame", False) or not rec.shape.has_text_frame:
        raise PatchError("%s (%s) has no text frame" % (rec.sid, rec.type))
    else:
        tf = rec.shape.text_frame
    write_paragraphs_inherit(tf, items)
    ctx.log.append("set-text slide %d %s (%d paragraph(s))" % (op["slide"], rec.sid, len(items)))


def op_swap_image(ctx, op):
    if op.get("media") and not op.get("rid") and not op.get("shape"):
        # global mode: overwrite the media part's bytes — EVERY reference
        # (slides, masters, layouts) changes; frames keep their geometry
        name = op["media"]
        image_path = Path(op["image"])
        parts = {
            str(p.partname).split("/")[-1]: p
            for p in ctx.prs.part.package.iter_parts()
            if str(p.partname).startswith("/ppt/media/")
        }
        if name not in parts:
            raise PatchError(
                "media '%s' not in package. Available media: %s" % (name, ", ".join(sorted(parts)))
            )
        parts[name]._blob = image_path.read_bytes()
        old_ext = name.rsplit(".", 1)[-1].lower()
        warn = ""
        if image_path.suffix.lstrip(".").lower() != old_ext:
            warn = " (note: .%s bytes under a .%s name — renders fine, same format is safer)" % (
                image_path.suffix.lstrip(".").lower(), old_ext)
        affected = []
        for i, s in enumerate(ctx.prs.slides):
            for rel in s.part.rels.values():
                if (rel.reltype == IMAGE_RELTYPE and not rel.is_external
                        and rel._target.partname == parts[name].partname):
                    affected.append(i)
                    break
        ctx.touched.update(affected)
        ctx.log.append(
            "swap-image media %s -> %s bytes; every reference changes (slides %s + any master/layout use)%s"
            % (name, image_path.name, ",".join(str(i) for i in affected) or "none directly", warn)
        )
        return
    slide_idx = op["slide"]
    slide = ctx.prs.slides[slide_idx]
    part = slide.part
    image_path = Path(op["image"])
    rid = op.get("rid")
    target_shape = None
    if op.get("shape"):
        rec = ctx.rec(slide_idx, op["shape"])
        if not rec.rids:
            raise PatchError("%s on slide %d carries no image" % (rec.sid, slide_idx))
        rid = rec.rids[0][0]
        target_shape = rec.shape
    if not rid:
        raise PatchError("swap-image needs 'rid' or 'shape'")
    if rid not in part.rels or part.rels[rid].reltype != IMAGE_RELTYPE:
        raise PatchError(
            "'%s' is not an image relationship on slide %d.\n%s"
            % (rid, slide_idx, slide_listing(ctx.index, slide_idx))
        )
    if target_shape is None:
        for r in ctx.index[slide_idx].values():
            if any(a == rid for a, _ in r.rids):
                target_shape = r.shape
                break
    blips = [
        b
        for b in slide._element.findall(".//" + qn("a:blip"))
        if b.get(R_EMBED) == rid
    ]
    _, new_rid = part.get_or_add_image_part(str(image_path))
    for b in blips:
        b.set(R_EMBED, new_rid)
    if not any(
        b.get(R_EMBED) == rid for b in slide._element.findall(".//" + qn("a:blip"))
    ):
        try:
            part.drop_rel(rid)
        except KeyError:
            pass
    # fit: auto = contain when aspect clearly differs from the frame
    fit = op.get("fit", "auto")
    if target_shape is not None and fit != "stretch":
        from PIL import Image

        with Image.open(image_path) as im:
            iw, ih = im.size
        bw, bh = target_shape.width, target_shape.height
        if iw and ih and bw and bh:
            frame_ar, img_ar = bw / bh, iw / ih
            mismatch = abs(frame_ar - img_ar) / frame_ar
            if fit == "contain" or (fit == "auto" and mismatch > 0.08):
                scale = min(bw / iw, bh / ih)
                nw, nh = int(iw * scale), int(ih * scale)
                target_shape.left = target_shape.left + (bw - nw) // 2
                target_shape.top = target_shape.top + (bh - nh) // 2
                target_shape.width = nw
                target_shape.height = nh
                ctx.log.append(
                    "swap-image slide %d %s -> %s (contain, aspect mismatch %.0f%%)"
                    % (slide_idx, rid, image_path.name, mismatch * 100)
                )
                ctx.reindex_slide(slide_idx)
                return
    ctx.log.append("swap-image slide %d %s -> %s" % (slide_idx, rid, image_path.name))
    ctx.reindex_slide(slide_idx)


def op_replace_text(ctx, op):
    frm, to = op["from"], op["to"]
    scope = op.get("scope", "deck")
    elements = []  # (label, lxml element)
    if scope == "deck":
        elements = [("slide %d" % i, s._element) for i, s in enumerate(ctx.prs.slides)]
    elif scope == "slide":
        if "slide" not in op:
            raise PatchError("replace-text scope 'slide' needs a 'slide' index")
        elements = [("slide %d" % op["slide"], ctx.prs.slides[op["slide"]]._element)]
    elif scope == "master":
        for mi, master in enumerate(ctx.prs.slide_masters):
            elements.append(("master %d" % mi, master._element))
            for li, layout in enumerate(master.slide_layouts):
                elements.append(("master %d layout %d" % (mi, li), layout._element))
    else:
        raise PatchError("replace-text scope must be deck|master|slide, got '%s'" % scope)
    count = 0
    where = []
    for label, el in elements:
        for t in el.iter(qn("a:t")):
            if t.text and frm in t.text:
                t.text = t.text.replace(frm, to)
                count += 1
                where.append(label)
    if count == 0:
        # help the agent: was it split across runs?
        hints = []
        for label, el in elements:
            for para in el.iter(qn("a:p")):
                joined = "".join(t.text or "" for t in para.iter(qn("a:t")))
                if frm in joined:
                    hints.append(label)
        msg = "replace-text: '%s' not found in scope '%s'" % (frm, scope)
        if hints:
            msg += (
                " as a contiguous run, but it IS present (split across runs) in: %s. "
                "Use set-text on those shapes instead." % ", ".join(sorted(set(hints)))
            )
        raise PatchError(msg)
    if scope in ("deck", "slide"):
        for label in set(where):
            ctx.touched.add(int(label.split()[1]))
    ctx.log.append(
        "replace-text [%s] '%s' -> '%s' (%d occurrence(s) in %s)"
        % (scope, frm, to, count, ", ".join(sorted(set(where))))
    )


def _hex(v, what):
    v = str(v).lstrip("#").upper()
    if not re.fullmatch(r"[0-9A-F]{6}", v):
        raise PatchError("%s must be a 6-digit hex color, got '%s'" % (what, v))
    return v


def op_replace_color(ctx, op):
    """Swap one concrete color for another everywhere in scope — the re-theme
    primitive. Touches srgbClr values in fills, lines, gradients, text, and
    effects; theme-indexed colors (schemeClr) are untouched by design (remap
    those in the theme, not per use)."""
    frm = _hex(op["from"], "replace-color from")
    to = _hex(op["to"], "replace-color to")
    scope = op.get("scope", "deck")
    elements = []
    if scope == "deck":
        elements = [("slide %d" % i, s._element) for i, s in enumerate(ctx.prs.slides)]
    elif scope == "slide":
        if "slide" not in op:
            raise PatchError("replace-color scope 'slide' needs a 'slide' index")
        elements = [("slide %d" % op["slide"], ctx.prs.slides[op["slide"]]._element)]
    elif scope == "master":
        for mi, master in enumerate(ctx.prs.slide_masters):
            elements.append(("master %d" % mi, master._element))
            for li, layout in enumerate(master.slide_layouts):
                elements.append(("master %d layout %d" % (mi, li), layout._element))
    else:
        raise PatchError("replace-color scope must be deck|master|slide, got '%s'" % scope)
    count = 0
    where = []
    for label, el in elements:
        for node in el.iter(qn("a:srgbClr")):
            if (node.get("val") or "").upper() == frm:
                node.set("val", to)
                count += 1
                where.append(label)
    if count == 0:
        # help the agent: list what IS there so it can correct in one pass
        seen = {}
        for label, el in elements:
            for node in el.iter(qn("a:srgbClr")):
                v = (node.get("val") or "").upper()
                seen[v] = seen.get(v, 0) + 1
        top = ", ".join("%s (x%d)" % kv for kv in sorted(seen.items(), key=lambda kv: -kv[1])[:12])
        raise PatchError(
            "replace-color: %s not found in scope '%s'. Colors present: %s" % (frm, scope, top)
        )
    if scope in ("deck", "slide"):
        for label in set(where):
            ctx.touched.add(int(label.split()[1]))
    ctx.log.append(
        "replace-color [%s] %s -> %s (%d occurrence(s) across %s)"
        % (scope, frm, to, count, ", ".join(sorted(set(where))))
    )


def op_set_notes(ctx, op):
    slide = ctx.prs.slides[op["slide"]]
    slide.notes_slide.notes_text_frame.text = op["notes"]
    ctx.log.append("set-notes slide %d (%d chars)" % (op["slide"], len(op["notes"])))


PROPS_KEYS = ("title", "subject", "author", "keywords", "comments", "category", "last_modified_by")


def op_set_props(ctx, op):
    """Document metadata (the File > Info panel). Strings only; set "" to clear."""
    cp = ctx.prs.core_properties
    changed = []
    for k in PROPS_KEYS:
        if k in op:
            setattr(cp, k, str(op[k]))
            changed.append(k)
    ctx.log.append("set-props %s" % ", ".join(changed))


# transition type -> (p: element local name, {option: valid values})
TRANSITION_TYPES = {
    "fade":     ("fade",     {}),
    "cut":      ("cut",      {}),
    "dissolve": ("dissolve", {}),
    "push":     ("push",     {"dir": ("l", "r", "u", "d")}),
    "wipe":     ("wipe",     {"dir": ("l", "r", "u", "d")}),
    "split":    ("split",    {"orient": ("horz", "vert"), "dir": ("in", "out")}),
    "cover":    ("cover",    {"dir": ("l", "r", "u", "d", "ld", "lu", "rd", "ru")}),
    "uncover":  ("pull",     {"dir": ("l", "r", "u", "d", "ld", "lu", "rd", "ru")}),
    "zoom":     ("zoom",     {"dir": ("in", "out")}),
}
TRANSITION_EL_TO_TYPE = {el: t for t, (el, _) in TRANSITION_TYPES.items()}
TRANSITION_SPEEDS = ("slow", "med", "fast")


def _transition_dict(slide):
    """Read a slide's p:transition into the same vocabulary set-slide takes."""
    el = slide._element.find(qn("p:transition"))
    if el is None:
        return None
    out = {}
    for child in el:
        local = child.tag.rsplit("}", 1)[-1]
        out["type"] = TRANSITION_EL_TO_TYPE.get(local, local)
        for k in ("dir", "orient"):
            if child.get(k):
                out[k] = child.get(k)
        break
    if el.get("spd"):
        out["speed"] = el.get("spd")
    if el.get("advTm") is not None:
        out["advance_after"] = int(el.get("advTm")) / 1000.0
    if el.get("advClick") == "0":
        out["advance_on_click"] = False
    return out


def op_set_slide(ctx, op):
    """Slide-level properties: hidden, background, transition."""
    slide = ctx.prs.slides[op["slide"]]
    done = []
    if "transition" in op:
        tr = op["transition"]
        sld = slide._element
        old = sld.find(qn("p:transition"))
        if old is not None:
            sld.remove(old)
        if tr == "none":
            done.append("transition removed")
        else:
            el_name, optspec = TRANSITION_TYPES[tr["type"]]
            tr_el = sld.makeelement(qn("p:transition"), {})
            if "speed" in tr:
                tr_el.set("spd", tr["speed"])
            if tr.get("advance_on_click") is False:
                tr_el.set("advClick", "0")
            if "advance_after" in tr:
                tr_el.set("advTm", str(int(round(float(tr["advance_after"]) * 1000))))
            eff = tr_el.makeelement(qn("p:" + el_name), {})
            for k in optspec:
                if k in tr:
                    eff.set(k, tr[k])
            tr_el.append(eff)
            # schema order: p:cSld, p:clrMapOvr, p:transition, p:timing
            anchor = sld.find(qn("p:clrMapOvr"))
            if anchor is None:
                anchor = sld.find(qn("p:cSld"))
            anchor.addnext(tr_el)
            done.append("transition=%s" % tr["type"])
    if "hidden" in op:
        if op["hidden"]:
            slide._element.set("show", "0")
        else:
            slide._element.attrib.pop("show", None)
        done.append("hidden=%s" % bool(op["hidden"]))
    if "background" in op:
        bg = op["background"]
        # clear any existing explicit background first
        csld = slide._element.find(qn("p:cSld"))
        old = csld.find(qn("p:bg"))
        if old is not None:
            csld.remove(old)
        if isinstance(bg, str):
            fill = slide.background.fill
            fill.solid()
            h = _hex(bg, "background")
            fill.fore_color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            done.append("background=%s" % h)
        elif isinstance(bg, dict) and "gradient" in bg:
            g = bg["gradient"]
            fill = slide.background.fill
            fill.gradient()
            stops = list(fill.gradient_stops)
            for stop, hexs in zip(stops, g["colors"]):
                h = _hex(hexs, "gradient color")
                stop.color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            if "angle" in g:
                fill.gradient_angle = float(g["angle"])
            done.append("background=gradient")
        elif isinstance(bg, dict) and "image" in bg:
            path = Path(bg["image"])
            if not path.exists():
                raise PatchError("background image not found: %s" % path)
            image_part, rid = slide.part.get_or_add_image_part(str(path))
            a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            bg_el = csld.makeelement(qn("p:bg"), {})
            bgpr = bg_el.makeelement(qn("p:bgPr"), {})
            blipfill = bgpr.makeelement("{%s}blipFill" % a, {"rotWithShape": "1"})
            blip = blipfill.makeelement("{%s}blip" % a, {})
            blip.set(qn("r:embed"), rid)
            stretch = blipfill.makeelement("{%s}stretch" % a, {})
            stretch.append(stretch.makeelement("{%s}fillRect" % a, {}))
            blipfill.append(blip)
            blipfill.append(stretch)
            bgpr.append(blipfill)
            bgpr.append(bgpr.makeelement("{%s}effectLst" % a, {}))
            bg_el.append(bgpr)
            csld.insert(0, bg_el)  # p:bg must be cSld's first child
            done.append("background=image %s" % path.name)
        else:
            raise PatchError(
                'set-slide "background" takes "RRGGBB", {"gradient": {"colors": [..], "angle": deg}}, '
                'or {"image": "/abs/path.png"}'
            )
    if not done:
        raise PatchError('set-slide: give "hidden", "background" and/or "transition"')
    ctx.touched.add(op["slide"])
    ctx.log.append("set-slide %d: %s" % (op["slide"], ", ".join(done)))


THEME_COLOR_KEYS = ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2", "accent3",
                    "accent4", "accent5", "accent6", "hlink", "folHlink")
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def op_set_theme(ctx, op):
    """Remap theme scheme colors and/or major/minor fonts. This is where
    template-driven decks keep their palette — replace-color handles literal
    colors; this handles the theme-indexed ones."""
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from lxml import etree

    masters = list(ctx.prs.slide_masters)
    if "master" in op:
        masters = [masters[op["master"]]]
    colors = op.get("colors", {})
    fonts = op.get("fonts", {})
    if not colors and not fonts:
        raise PatchError('set-theme: give "colors" and/or "fonts"')
    done = []
    for mi, master in enumerate(masters):
        tp = master.part.part_related_by(RT.THEME)
        editable = hasattr(tp, "_element")
        root = tp._element if editable else etree.fromstring(tp.blob)
        nsmap = {"a": A_NS}
        for key, hexv in colors.items():
            el = root.find("a:themeElements/a:clrScheme/a:%s" % key, nsmap)
            if el is None:
                raise PatchError("theme color slot '%s' not found in master %d's theme" % (key, mi))
            h = _hex(hexv, "theme color %s" % key)
            for ch in list(el):
                el.remove(ch)
            s = etree.SubElement(el, "{%s}srgbClr" % A_NS)
            s.set("val", h)
        for role, name in fonts.items():
            tag = {"major": "a:majorFont", "minor": "a:minorFont"}.get(role)
            if tag is None:
                raise PatchError('set-theme fonts keys are "major" and "minor", got \'%s\'' % role)
            latin = root.find("a:themeElements/a:fontScheme/%s/a:latin" % tag, nsmap)
            if latin is None:
                raise PatchError("theme %s font slot not found in master %d's theme" % (role, mi))
            latin.set("typeface", str(name))
        if not editable:
            tp._blob = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
        done.append("master %d" % mi)
    ctx.log.append(
        "set-theme %s: %s%s" % (
            ", ".join(done),
            ("colors " + ", ".join("%s->%s" % (k, v) for k, v in colors.items())) if colors else "",
            ((" " if colors else "") + "fonts " + ", ".join("%s=%s" % (k, v) for k, v in fonts.items())) if fonts else "",
        )
    )


def op_move(ctx, op):
    rec = ctx.rec(op["slide"], op["shape"])
    sh = rec.shape
    if "to" in op:
        ax, ay = (Inches(v) for v in op["to"])
        sh.left = int(rec.tf.child_x(ax))
        sh.top = int(rec.tf.child_y(ay))
    elif "by" in op:
        dx, dy = (Inches(v) for v in op["by"])
        sh.left = (sh.left or 0) + int(dx / rec.tf.sx)
        sh.top = (sh.top or 0) + int(dy / rec.tf.sy)
    else:
        raise PatchError("move needs 'to':[l,t] or 'by':[dx,dy] (inches)")
    ctx.reindex_slide(op["slide"])
    ctx.log.append("move slide %d %s" % (op["slide"], rec.sid))


def op_resize(ctx, op):
    rec = ctx.rec(op["slide"], op["shape"])
    sh = rec.shape
    if "size" in op:
        w, h = op["size"]
        sh.width = int(Inches(w) / rec.tf.sx)
        sh.height = int(Inches(h) / rec.tf.sy)
    elif "scale" in op:
        f = float(op["scale"])
        sh.width = int((sh.width or 0) * f)
        sh.height = int((sh.height or 0) * f)
    else:
        raise PatchError("resize needs 'size':[w,h] (inches) or 'scale':f")
    ctx.reindex_slide(op["slide"])
    ctx.log.append("resize slide %d %s" % (op["slide"], rec.sid))


def _iter_text_frames(rec):
    if rec.is_table:
        for row in rec.shape.table.rows:
            for cell in row.cells:
                yield cell.text_frame
    elif getattr(rec.shape, "has_text_frame", False) and rec.shape.has_text_frame:
        yield rec.shape.text_frame


def _apply_style_keys(rec, op):
    """Shared by set-style and add-shape: shape-wide fonts, fill/gradient,
    line/border, rotation. Returns how many things it changed."""
    styled = 0
    for tf in _iter_text_frames(rec):
        for para in tf.paragraphs:
            for run in para.runs:
                if "font_size" in op:
                    run.font.size = Pt(op["font_size"])
                if "font_name" in op:
                    run.font.name = op["font_name"]
                if "bold" in op:
                    run.font.bold = op["bold"]
                if "italic" in op:
                    run.font.italic = op["italic"]
                if "underline" in op:
                    run.font.underline = op["underline"]
                if "color" in op:
                    h = op["color"].lstrip("#")
                    run.font.color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
                styled += 1
    if "fill" in op:
        try:
            if op["fill"] == "none":
                rec.shape.fill.background()
            else:
                h = op["fill"].lstrip("#")
                rec.shape.fill.solid()
                rec.shape.fill.fore_color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            styled += 1
        except Exception as e:
            raise PatchError("cannot set fill on %s (%s): %s" % (rec.sid, rec.type, e))
    if "gradient" in op:
        g = op["gradient"]
        try:
            fill = rec.shape.fill
            fill.solid()      # reset so gradient() always yields a fresh 2-stop fill
            fill.gradient()
            stops = list(fill.gradient_stops)
            for stop, hexs in zip(stops, g["colors"]):
                h = hexs.lstrip("#")
                stop.color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            for stop, pos in zip(stops, g.get("positions", [])):
                stop.position = float(pos)
            if "angle" in g:
                fill.gradient_angle = float(g["angle"])
            styled += 1
        except Exception as e:
            raise PatchError("cannot set gradient on %s (%s): %s" % (rec.sid, rec.type, e))
    if any(k in op for k in ("line", "line_color", "line_width", "line_dash")):
        try:
            ln = rec.shape.line
            if op.get("line") == "none":
                ln.fill.background()
            if "line_color" in op:
                h = op["line_color"].lstrip("#")
                ln.color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            if "line_width" in op:
                ln.width = Pt(float(op["line_width"]))
            if "line_dash" in op:
                ln.dash_style = LINE_DASH_STYLES[op["line_dash"]]
            styled += 1
        except PatchError:
            raise
        except Exception as e:
            raise PatchError("cannot set line on %s (%s): %s" % (rec.sid, rec.type, e))
    if "rotation" in op:
        try:
            rec.shape.rotation = float(op["rotation"])
            styled += 1
        except Exception as e:
            raise PatchError("cannot set rotation on %s (%s): %s" % (rec.sid, rec.type, e))
    if "insets" in op:
        ins = op["insets"]
        if not (isinstance(ins, (list, tuple)) and len(ins) == 4):
            raise PatchError('"insets" must be [left, top, right, bottom] in inches')
        if not getattr(rec.shape, "has_text_frame", False):
            raise PatchError("cannot set insets on %s (%s): no text frame" % (rec.sid, rec.type))
        tf = rec.shape.text_frame
        tf.margin_left, tf.margin_top = Inches(ins[0]), Inches(ins[1])
        tf.margin_right, tf.margin_bottom = Inches(ins[2]), Inches(ins[3])
        styled += 1
    if "adjustments" in op:
        try:
            adjs = rec.shape.adjustments
            for i, val in enumerate(op["adjustments"]):
                adjs[i] = float(val)
            styled += 1
        except Exception as e:
            raise PatchError("cannot set adjustments on %s (%s): %s" % (rec.sid, rec.type, e))
    if "shadow" in op:
        try:
            rec.shape.shadow.inherit = bool(op["shadow"])
            if not op["shadow"]:
                # an empty effectLst SHOULD override the style's effectRef, but
                # some renderers (LibreOffice) still apply it — zero it too
                ref = rec.shape._element.find(qn("p:style") + "/" + qn("a:effectRef"))
                if ref is not None:
                    ref.set("idx", "0")
            styled += 1
        except Exception as e:
            raise PatchError("cannot set shadow on %s (%s): %s" % (rec.sid, rec.type, e))
    if "alt_text" in op:
        _set_alt_text(rec.shape, op["alt_text"])
        styled += 1
    return styled


def _set_alt_text(shape, text):
    # the shape's own cNvPr is its first descendant cNvPr in document order
    cnvpr = shape._element.find(".//" + qn("p:cNvPr"))
    if cnvpr is None:
        raise PatchError("shape has no cNvPr element (cannot set alt_text)")
    if text:
        cnvpr.set("descr", str(text))
    else:
        cnvpr.attrib.pop("descr", None)


def _get_alt_text(shape):
    cnvpr = shape._element.find(".//" + qn("p:cNvPr"))
    return cnvpr.get("descr") if cnvpr is not None else None


def op_set_style(ctx, op):
    rec = ctx.rec(op["slide"], op["shape"])
    styled = _apply_style_keys(rec, op)
    if styled == 0:
        raise PatchError("set-style touched nothing on %s (no runs/fill/line?)" % rec.sid)
    ctx.log.append("set-style slide %d %s (%d run(s)/fill/line)" % (op["slide"], rec.sid, styled))


def _prune_timing_refs(slide_el, native_id):
    """Drop animation effects and build entries that target a shape id, so a
    deleted shape never leaves dangling spTgt references in the timing tree."""
    timing = slide_el.find(qn("p:timing"))
    if timing is None:
        return
    for tgt in list(timing.findall(".//" + qn("p:spTgt"))):
        if tgt.get("spid") != str(native_id):
            continue
        node = tgt  # remove the innermost p:par (the effect's own timing node)
        while node is not None and node.tag != qn("p:par"):
            node = node.getparent()
        if node is not None and node.getparent() is not None:
            node.getparent().remove(node)
    for bld in list(timing.findall(".//" + qn("p:bldP"))):
        if bld.get("spid") == str(native_id):
            bld.getparent().remove(bld)
    changed = True
    while changed:  # cascade: drop containers the removal left empty
        changed = False
        for ctl in timing.findall(".//" + qn("p:childTnLst")):
            if len(ctl) == 0:
                holder = ctl.getparent().getparent()  # the par/seq around the cTn
                if holder is not None and holder.getparent() is not None:
                    holder.getparent().remove(holder)
                    changed = True
                    break
        if not changed:
            for lst in (timing.find(qn("p:bldLst")), timing.find(qn("p:tnLst"))):
                if lst is not None and len(lst) == 0:
                    timing.remove(lst)
                    changed = True
    if len(timing) == 0:
        slide_el.remove(timing)


def op_delete(ctx, op):
    rec = ctx.rec(op["slide"], op["shape"])
    el = rec.shape._element
    slide_el = ctx.prs.slides[op["slide"]]._element
    ids = [c.get("id") for c in el.iter(qn("p:cNvPr")) if c.get("id")]
    el.getparent().remove(el)
    for native_id in ids:  # a group brings its members' animations with it
        _prune_timing_refs(slide_el, native_id)
    ctx.reindex_slide(op["slide"])
    ctx.log.append("delete slide %d %s (%s)" % (op["slide"], rec.sid, rec.name))


def _next_shape_id(slide_part_element):
    ids = [
        int(e.get("id"))
        for e in slide_part_element.iter()
        if e.tag.endswith("}cNvPr") and e.get("id") and e.get("id").isdigit()
    ]
    return (max(ids) + 1) if ids else 1


def _reassign_ids(new_el, start_id):
    """Give every cNvPr in a copied subtree a fresh unique id. Returns the id
    of the subtree's ROOT shape and the next free id."""
    next_id = start_id
    root_id = None
    for e in new_el.iter():
        if e.tag.endswith("}cNvPr"):
            if root_id is None:
                root_id = next_id
            e.set("id", str(next_id))
            name = e.get("name") or "Shape"
            e.set("name", name + " copy")
            next_id += 1
    return root_id, next_id


def _find_shape_by_id(slide, shape_id):
    def walk(shapes):
        for sh in shapes:
            if sh.shape_id == shape_id:
                return sh
            try:
                if sh.shape_type == GROUP_TYPE:
                    found = walk(sh.shapes)
                    if found is not None:
                        return found
            except Exception:
                pass
        return None

    return walk(slide.shapes)


def op_duplicate(ctx, op):
    slide_idx = op["slide"]
    rec = ctx.rec(slide_idx, op["shape"])
    slide = ctx.prs.slides[slide_idx]
    src_el = rec.shape._element
    new_el = copy.deepcopy(src_el)
    root_id, _ = _reassign_ids(new_el, _next_shape_id(slide._element))
    src_el.addnext(new_el)
    ctx.reindex_slide(slide_idx)
    new_sid = "s%d" % root_id
    new_rec = ctx.index[slide_idx][new_sid]
    if "offset" in op:
        dx, dy = op["offset"]
        sh = new_rec.shape
        sh.left = (sh.left or 0) + int(Inches(dx) / new_rec.tf.sx)
        sh.top = (sh.top or 0) + int(Inches(dy) / new_rec.tf.sy)
    if "at" in op:
        ax, ay = op["at"]
        sh = new_rec.shape
        sh.left = int(new_rec.tf.child_x(Inches(ax)))
        sh.top = int(new_rec.tf.child_y(Inches(ay)))
    if "text" in op:
        if not getattr(new_rec.shape, "has_text_frame", False) or not new_rec.shape.has_text_frame:
            raise PatchError("duplicate: new shape %s has no text frame for 'text'" % new_sid)
        write_paragraphs_inherit(new_rec.shape.text_frame, norm_paragraph_items(op["text"]))
    ctx.reindex_slide(slide_idx)
    ctx.log.append(
        "duplicate slide %d %s -> %s%s"
        % (slide_idx, rec.sid, new_sid, (" (text set)" if "text" in op else ""))
    )


def op_copy_shape(ctx, op):
    src_idx, dst_idx = op["from_slide"], op["slide"]
    rec = ctx.rec(src_idx, op["shape"])
    src_part = ctx.prs.slides[src_idx].part
    dst_slide = ctx.prs.slides[dst_idx]
    dst_part = dst_slide.part
    new_el = copy.deepcopy(rec.shape._element)
    root_id, _ = _reassign_ids(new_el, _next_shape_id(dst_slide._element))
    # re-home relationship references (images, hyperlinks)
    for e in new_el.iter():
        for attr in (R_EMBED, R_LINK, R_ID):
            rid = e.get(attr)
            if not rid:
                continue
            if rid not in src_part.rels:
                continue
            rel = src_part.rels[rid]
            if rel.is_external:
                new_rid = dst_part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_rid = dst_part.rels.get_or_add(rel.reltype, rel._target)
            e.set(attr, new_rid)
    dst_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    ctx.reindex_slide(dst_idx)
    new_sid = "s%d" % root_id
    new_rec = ctx.index[dst_idx][new_sid]
    if "at" in op:
        ax, ay = op["at"]
        new_rec.shape.left = int(Inches(ax))
        new_rec.shape.top = int(Inches(ay))
    if "text" in op:
        write_paragraphs_inherit(new_rec.shape.text_frame, norm_paragraph_items(op["text"]))
    ctx.reindex_slide(dst_idx)
    ctx.log.append("copy-shape slide %d %s -> slide %d %s" % (src_idx, rec.sid, dst_idx, new_sid))


def op_add_shape(ctx, op):
    slide_idx = op["slide"]
    slide = ctx.prs.slides[slide_idx]
    kind = resolve_shape_kind(op.get("kind", "textbox"))
    if kind == "line":
        if "from" not in op or "to" not in op:
            raise PatchError('add-shape kind "line" needs "from":[x,y] and "to":[x,y] (inches)')
        fx, fy = op["from"]
        tx, ty = op["to"]
        sh = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(fx), Inches(fy), Inches(tx), Inches(ty)
        )
    else:
        if "at" not in op or "size" not in op:
            raise PatchError('add-shape needs "at":[l,t] and "size":[w,h] (inches)')
        l, t = op["at"]
        w, h = op["size"]
        if kind == "textbox":
            sh = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            sh.text_frame.word_wrap = True
        else:
            sh = slide.shapes.add_shape(kind, Inches(l), Inches(t), Inches(w), Inches(h))
    if "name" in op:
        sh.name = op["name"]
    ctx.reindex_slide(slide_idx)
    sid = "s%d" % sh.shape_id
    rec = ctx.index[slide_idx][sid]
    if "text" in op:
        if not getattr(sh, "has_text_frame", False) or not sh.has_text_frame:
            raise PatchError('add-shape kind "%s" cannot carry "text"' % op.get("kind"))
        write_paragraphs_inherit(sh.text_frame, norm_paragraph_items(op["text"]))
    _apply_style_keys(rec, op)
    ctx.reindex_slide(slide_idx)
    ctx.log.append("add-shape slide %d %s -> %s" % (slide_idx, op.get("kind", "textbox"), sid))


def op_add_picture(ctx, op):
    slide_idx = op["slide"]
    slide = ctx.prs.slides[slide_idx]
    l, t = op["at"]
    kw = {}
    if "size" in op:
        kw = {"width": Inches(op["size"][0]), "height": Inches(op["size"][1])}
    elif "width" in op:
        kw = {"width": Inches(op["width"])}  # height keeps aspect ratio
    elif "height" in op:
        kw = {"height": Inches(op["height"])}
    sh = slide.shapes.add_picture(str(op["image"]), Inches(l), Inches(t), **kw)
    if "crop" in op:
        c = op["crop"]
        if not (isinstance(c, (list, tuple)) and len(c) == 4):
            raise PatchError('"crop" must be [left, top, right, bottom] fractions (0-1)')
        sh.crop_left, sh.crop_top, sh.crop_right, sh.crop_bottom = [float(v) for v in c]
    if "shadow" in op:
        sh.shadow.inherit = bool(op["shadow"])
    if "alt_text" in op:
        _set_alt_text(sh, op["alt_text"])
    ctx.reindex_slide(slide_idx)
    ctx.log.append(
        "add-picture slide %d %s -> s%d (%.2fx%.2fin)"
        % (slide_idx, Path(op["image"]).name, sh.shape_id, inches(sh.width), inches(sh.height))
    )


def op_add_table(ctx, op):
    slide_idx = op["slide"]
    rows = op["rows"]
    l, t = op["at"]
    w, h = op["size"]
    gf = ctx.prs.slides[slide_idx].shapes.add_table(
        len(rows), len(rows[0]), Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if "name" in op:
        gf.name = op["name"]
    tbl = gf.table
    if "first_row" in op:
        tbl.first_row = bool(op["first_row"])
    if "banding" in op:
        tbl.horz_banding = bool(op["banding"])
    if "col_widths" in op:
        widths = op["col_widths"]
        if len(widths) != len(rows[0]):
            raise PatchError(
                "col_widths has %d entries for %d columns" % (len(widths), len(rows[0]))
            )
        for col, cw in zip(tbl.columns, widths):
            col.width = Inches(cw)
    fill = op.get("fill")
    fills = op.get("fills")  # row-major grid; entries: "RRGGBB" | "none" | null
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            cf = fills[ri][ci] if fills else fill
            if cf == "none":
                cell.fill.background()
            elif cf:
                h = cf.lstrip("#")
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            if "font_size" in op or "color" in op:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if "font_size" in op:
                            run.font.size = Pt(op["font_size"])
                        if "color" in op:
                            ch = op["color"].lstrip("#")
                            run.font.color.rgb = RGBColor(
                                int(ch[0:2], 16), int(ch[2:4], 16), int(ch[4:6], 16)
                            )
    ctx.reindex_slide(slide_idx)
    ctx.log.append(
        "add-table slide %d -> s%d (%dx%d)" % (slide_idx, gf.shape_id, len(rows), len(rows[0]))
    )


def _resolve_layout(prs, ref):
    layouts = [l for m in prs.slide_masters for l in m.slide_layouts]
    if ref is None:  # prefer a blank layout, else the emptiest
        blank = [l for l in layouts if "blank" in l.name.lower()]
        return blank[0] if blank else min(layouts, key=lambda l: len(l.shapes))
    if isinstance(ref, int):
        if not (0 <= ref < len(layouts)):
            raise PatchError("layout %d out of range (0-%d)" % (ref, len(layouts) - 1))
        return layouts[ref]
    named = [l for l in layouts if l.name == ref]
    if not named:
        named = [l for l in layouts if str(ref).lower() in l.name.lower()]
    if len(named) != 1:
        raise PatchError(
            "layout '%s' %s — available: %s"
            % (ref, "not found" if not named else "ambiguous",
               ", ".join("'%s'" % l.name for l in layouts))
        )
    return named[0]


def op_add_slide(ctx, op):
    prs = ctx.prs
    layout = _resolve_layout(prs, op.get("layout"))
    prs.slides.add_slide(layout)
    new_idx = len(prs.slides._sldIdLst) - 1
    at = op.get("at")
    if at is not None and at != new_idx:
        if not (0 <= at <= new_idx):
            raise PatchError("add-slide 'at' %s out of range (0-%d)" % (at, new_idx))
        lst = prs.slides._sldIdLst
        el = lst[new_idx]
        lst.remove(el)
        lst.insert(at, el)
        ctx.touched = {(s + 1 if s >= at else s) for s in ctx.touched}
        new_idx = at
    ctx.touched.add(new_idx)
    ctx.index = build_index(prs, measure=False)  # indices shifted — full rebuild
    ctx.log.append("add-slide layout '%s' -> slide %d" % (layout.name, new_idx))


SHAPE_TAGS = {qn("p:sp"), qn("p:grpSp"), qn("p:pic"), qn("p:graphicFrame"), qn("p:cxnSp")}


def op_reorder(ctx, op):
    rec = ctx.rec(op["slide"], op["shape"])
    el = rec.shape._element
    parent = el.getparent()
    sibs = [c for c in parent if c.tag in SHAPE_TAGS]
    i = sibs.index(el)
    z = op.get("z")
    if z == "front":
        if sibs[-1] is not el:
            sibs[-1].addnext(el)
    elif z == "back":
        if sibs[0] is not el:
            sibs[0].addprevious(el)
    elif z == "forward":
        if i < len(sibs) - 1:
            sibs[i + 1].addnext(el)
    elif z == "backward":
        if i > 0:
            sibs[i - 1].addprevious(el)
    else:
        raise PatchError('reorder needs "z": front|back|forward|backward')
    ctx.reindex_slide(op["slide"])
    ctx.log.append("reorder slide %d %s -> %s" % (op["slide"], rec.sid, z))


A_TR, A_TC, A_GRIDCOL = qn("a:tr"), qn("a:tc"), qn("a:gridCol")
MERGE_ATTRS = ("gridSpan", "rowSpan", "hMerge", "vMerge")


def _table_tbl(ctx, op):
    rec = ctx.rec(op["slide"], op["shape"])
    if not rec.is_table:
        raise PatchError("%s on slide %d is not a table" % (rec.sid, op["slide"]))
    tbl = rec.shape.table._tbl
    for tc in tbl.iter(A_TC):
        if any(tc.get(a) for a in MERGE_ATTRS):
            raise PatchError(
                "%s has merged cells — row/col ops would corrupt the merge; use the xml escape hatch"
                % rec.sid
            )
    return rec, tbl


def _norm_idx(idx, n, what):
    if idx is None:
        return n - 1
    if not (-n <= idx < n):
        raise PatchError("%s %d out of range — table has %d" % (what, idx, n))
    return idx % n


def op_add_row(ctx, op):
    rec, tbl = _table_tbl(ctx, op)
    trs = tbl.findall(A_TR)
    src = trs[_norm_idx(op.get("copy"), len(trs), "copy row")]
    new_tr = copy.deepcopy(src)
    at = op.get("at")
    if at is None or at >= len(trs):
        trs[-1].addnext(new_tr)
        at = len(trs)
    else:
        trs[_norm_idx(at, len(trs), "row")].addprevious(new_tr)
        at = _norm_idx(at, len(trs), "row")
    table = rec.shape.table
    ncols = len(table.columns)
    cells = op.get("cells", [""] * ncols)
    if len(cells) != ncols:
        raise PatchError("'cells' has %d value(s) but the table has %d column(s)" % (len(cells), ncols))
    for ci, val in enumerate(cells):  # inherits the copied row's formatting
        write_paragraphs_inherit(table.cell(at, ci).text_frame, norm_paragraph_items(val))
    ctx.log.append("add-row slide %d %s at row %d" % (op["slide"], rec.sid, at))


def op_delete_row(ctx, op):
    rec, tbl = _table_tbl(ctx, op)
    trs = tbl.findall(A_TR)
    if len(trs) == 1:
        raise PatchError("cannot delete the last remaining row of %s" % rec.sid)
    idx = _norm_idx(op.get("row"), len(trs), "row")
    tbl.remove(trs[idx])
    ctx.log.append("delete-row slide %d %s row %d" % (op["slide"], rec.sid, idx))


def _rescale_grid(grid, target_total):
    cols = grid.findall(A_GRIDCOL)
    total = sum(int(c.get("w")) for c in cols)
    for c in cols:
        c.set("w", str(max(1, int(int(c.get("w")) * target_total / total))))


def op_add_col(ctx, op):
    rec, tbl = _table_tbl(ctx, op)
    grid = tbl.find(qn("a:tblGrid"))
    cols = grid.findall(A_GRIDCOL)
    ci = _norm_idx(op.get("copy"), len(cols), "copy col")
    total_w = sum(int(c.get("w")) for c in cols)
    at = op.get("at")
    append = at is None or at >= len(cols)
    new_col = copy.deepcopy(cols[ci])
    if append:
        cols[-1].addnext(new_col)
        at = len(cols)
    else:
        at = _norm_idx(at, len(cols), "col")
        cols[at].addprevious(new_col)
    _rescale_grid(grid, total_w)  # keep the table's overall width
    for tr in tbl.findall(A_TR):
        tcs = tr.findall(A_TC)
        new_tc = copy.deepcopy(tcs[ci])
        if append:
            tcs[-1].addnext(new_tc)
        else:
            tcs[at].addprevious(new_tc)
    table = rec.shape.table
    nrows = len(table.rows)
    cells = op.get("cells", [""] * nrows)
    if len(cells) != nrows:
        raise PatchError("'cells' has %d value(s) but the table has %d row(s)" % (len(cells), nrows))
    for ri, val in enumerate(cells):
        write_paragraphs_inherit(table.cell(ri, at).text_frame, norm_paragraph_items(val))
    ctx.log.append("add-col slide %d %s at col %d" % (op["slide"], rec.sid, at))


def op_delete_col(ctx, op):
    rec, tbl = _table_tbl(ctx, op)
    grid = tbl.find(qn("a:tblGrid"))
    cols = grid.findall(A_GRIDCOL)
    if len(cols) == 1:
        raise PatchError("cannot delete the last remaining column of %s" % rec.sid)
    idx = _norm_idx(op.get("col"), len(cols), "col")
    total_w = sum(int(c.get("w")) for c in cols)
    grid.remove(cols[idx])
    _rescale_grid(grid, total_w)
    for tr in tbl.findall(A_TR):
        tcs = tr.findall(A_TC)
        tr.remove(tcs[idx])
    ctx.log.append("delete-col slide %d %s col %d" % (op["slide"], rec.sid, idx))


OP_HANDLERS = {
    "set-text": op_set_text,
    "swap-image": op_swap_image,
    "replace-text": op_replace_text,
    "replace-color": op_replace_color,
    "set-notes": op_set_notes,
    "set-props": op_set_props,
    "set-slide": op_set_slide,
    "set-theme": op_set_theme,
    "move": op_move,
    "resize": op_resize,
    "set-style": op_set_style,
    "delete": op_delete,
    "duplicate": op_duplicate,
    "copy-shape": op_copy_shape,
    "add-shape": op_add_shape,
    "add-picture": op_add_picture,
    "add-table": op_add_table,
    "add-slide": op_add_slide,
    "reorder": op_reorder,
    "add-row": op_add_row,
    "delete-row": op_delete_row,
    "add-col": op_add_col,
    "delete-col": op_delete_col,
}


def validate_ops(ctx, ops):
    """Pre-flight every op; collect ALL errors before touching anything."""
    errors = []
    n_orig = len(ctx.prs.slides._sldIdLst)
    n_slides = n_orig  # virtual count — add-slide ops grow it as we scan
    created = set()  # (slide, name) pairs earlier add-shape ops will create
    for i, op in enumerate(ops):
        tag = "op[%d] %s" % (i, op.get("op", "?"))
        kind = op.get("op")
        if kind not in OP_HANDLERS:
            errors.append("%s: unknown op (valid: %s)" % (tag, ", ".join(VALID_OPS)))
            continue
        if kind == "add-slide":
            at = op.get("at")
            if at is not None and not (0 <= at <= n_slides):
                errors.append("%s: 'at' %s out of range (0-%d)" % (tag, at, n_slides))
            try:
                _resolve_layout(ctx.prs, op.get("layout"))
            except PatchError as e:
                errors.append("%s: %s" % (tag, e))
            n_slides += 1
            continue
        # slide bounds
        for key in ("slide", "from_slide"):
            if key in op and not (0 <= op[key] < n_slides):
                errors.append("%s: %s %s out of range (0-%d) — indices are 0-BASED" % (tag, key, op[key], n_slides - 1))
        if any(key in op and not (0 <= op[key] < n_slides) for key in ("slide", "from_slide")):
            continue
        if any(key in op and op[key] >= n_orig for key in ("slide", "from_slide")):
            continue  # targets a slide an earlier add-slide creates — checkable only at run time
        if kind in ("add-shape", "add-table") and "name" in op:
            created.add((op["slide"], op["name"]))
        if kind in ("set-text", "move", "resize", "set-style", "delete", "duplicate", "reorder",
                    "add-row", "delete-row", "add-col", "delete-col"):
            if (op.get("slide"), op.get("shape")) not in created:
                try:
                    ctx.rec(op["slide"], op["shape"])
                except KeyError:
                    errors.append("%s: missing required field" % tag)
                except PatchError as e:
                    errors.append("%s: %s" % (tag, e))
        if kind == "copy-shape":
            try:
                ctx.rec(op["from_slide"], op["shape"])
            except PatchError as e:
                errors.append("%s: %s" % (tag, e))
        if kind in ("swap-image", "add-picture"):
            if kind == "swap-image" and not op.get("rid") and not op.get("shape") and not op.get("media"):
                errors.append("%s: needs 'rid', 'shape', or 'media' (global byte swap)" % tag)
            if kind == "swap-image" and op.get("media") is None and "slide" not in op:
                errors.append("%s: needs 'slide' (or use 'media' for a deck-wide swap)" % tag)
            if kind == "add-picture" and "at" not in op:
                errors.append("%s: needs 'at':[l,t] (inches)" % tag)
            img = op.get("image")
            if not img or not Path(img).exists():
                errors.append("%s: image not found: %s" % (tag, img))
        if kind == "add-shape":
            try:
                k = resolve_shape_kind(op.get("kind", "textbox"))
                if k == "line":
                    if "from" not in op or "to" not in op:
                        errors.append('%s: kind "line" needs "from":[x,y] and "to":[x,y]' % tag)
                elif "at" not in op or "size" not in op:
                    errors.append('%s: needs "at":[l,t] and "size":[w,h] (inches)' % tag)
            except PatchError as e:
                errors.append("%s: %s" % (tag, e))
            if "text" in op:
                try:
                    norm_paragraph_items(op["text"])
                except PatchError as e:
                    errors.append("%s: %s" % (tag, e))
        if kind == "add-table":
            rows = op.get("rows")
            if not isinstance(rows, list) or not rows or not all(
                isinstance(r, list) and len(r) == len(rows[0]) and r for r in rows
            ):
                errors.append("%s: 'rows' must be a non-empty rectangular array of arrays" % tag)
            if "at" not in op or "size" not in op:
                errors.append('%s: needs "at":[l,t] and "size":[w,h] (inches)' % tag)
        if kind == "reorder" and op.get("z") not in ("front", "back", "forward", "backward"):
            errors.append('%s: needs "z": front|back|forward|backward' % tag)
        if kind == "replace-text" and ("from" not in op or "to" not in op):
            errors.append("%s: needs 'from' and 'to'" % tag)
        if kind == "replace-color":
            if "from" not in op or "to" not in op:
                errors.append("%s: needs 'from' and 'to' (6-digit hex)" % tag)
            else:
                for k in ("from", "to"):
                    try:
                        _hex(op[k], "'%s'" % k)
                    except PatchError as e:
                        errors.append("%s: %s" % (tag, e))
        if kind == "set-notes" and "notes" not in op:
            errors.append("%s: needs 'notes'" % tag)
        if kind == "set-props":
            given = [k for k in op if k != "op"]
            bad = sorted(k for k in given if k not in PROPS_KEYS)
            if bad:
                errors.append("%s: unknown key(s) %s — valid: %s" % (tag, ", ".join(bad), ", ".join(PROPS_KEYS)))
            if not [k for k in given if k in PROPS_KEYS]:
                errors.append("%s: give at least one of %s" % (tag, ", ".join(PROPS_KEYS)))
        if kind == "set-slide":
            if "slide" not in op:
                errors.append("%s: needs 'slide'" % tag)
            if not ("hidden" in op or "background" in op or "transition" in op):
                errors.append("%s: give \"hidden\", \"background\" and/or \"transition\"" % tag)
            if "transition" in op and op["transition"] != "none":
                tr = op["transition"]
                if not isinstance(tr, dict) or "type" not in tr:
                    errors.append(
                        '%s: "transition" takes {"type":"fade",...} or "none" — types: %s'
                        % (tag, ", ".join(sorted(TRANSITION_TYPES)))
                    )
                elif tr["type"] not in TRANSITION_TYPES:
                    errors.append(
                        "%s: transition type '%s' unknown — types: %s (or \"none\" to remove)"
                        % (tag, tr["type"], ", ".join(sorted(TRANSITION_TYPES)))
                    )
                else:
                    _, optspec = TRANSITION_TYPES[tr["type"]]
                    for k, v in tr.items():
                        if k == "type":
                            continue
                        elif k == "speed":
                            if v not in TRANSITION_SPEEDS:
                                errors.append("%s: transition speed '%s' — valid: %s" % (tag, v, ", ".join(TRANSITION_SPEEDS)))
                        elif k == "advance_after":
                            if not isinstance(v, (int, float)) or isinstance(v, bool) or v < 0:
                                errors.append("%s: transition advance_after must be seconds >= 0" % tag)
                        elif k == "advance_on_click":
                            if not isinstance(v, bool):
                                errors.append("%s: transition advance_on_click must be true/false" % tag)
                        elif k in optspec:
                            if v not in optspec[k]:
                                errors.append(
                                    "%s: transition %s '%s' invalid for %s — valid: %s"
                                    % (tag, k, v, tr["type"], ", ".join(optspec[k]))
                                )
                        else:
                            valid = ["type", "speed", "advance_after", "advance_on_click"] + sorted(optspec)
                            errors.append(
                                "%s: transition key '%s' not valid for %s — valid: %s"
                                % (tag, k, tr["type"], ", ".join(valid))
                            )
        if kind == "set-theme":
            if not (isinstance(op.get("colors"), dict) or isinstance(op.get("fonts"), dict)):
                errors.append("%s: give \"colors\" (dict) and/or \"fonts\" (dict)" % tag)
            n_masters = len(ctx.prs.slide_masters._sldMasterIdLst) if hasattr(ctx.prs.slide_masters, "_sldMasterIdLst") else len(list(ctx.prs.slide_masters))
            if "master" in op and not (0 <= op["master"] < n_masters):
                errors.append("%s: master %s out of range (0-%d)" % (tag, op["master"], n_masters - 1))
            for key in (op.get("colors") or {}):
                if key not in THEME_COLOR_KEYS:
                    errors.append("%s: unknown color slot '%s' — slots: %s" % (tag, key, ", ".join(THEME_COLOR_KEYS)))
            for key in (op.get("fonts") or {}):
                if key not in ("major", "minor"):
                    errors.append("%s: fonts keys are \"major\"/\"minor\", got '%s'" % (tag, key))
        if kind == "set-text":
            if "text" not in op:
                errors.append("%s: needs 'text'" % tag)
            else:
                try:
                    norm_paragraph_items(op["text"])
                except PatchError as e:
                    errors.append("%s: %s" % (tag, e))
        if kind in ("set-style", "add-shape"):
            if "line_dash" in op and op["line_dash"] not in LINE_DASH_STYLES:
                errors.append(
                    "%s: line_dash '%s' invalid (valid: %s)"
                    % (tag, op["line_dash"], ", ".join(sorted(LINE_DASH_STYLES)))
                )
            if "line" in op and op["line"] != "none":
                errors.append("%s: 'line' only accepts \"none\" — use line_color/line_width/line_dash to style it" % tag)
            if "gradient" in op:
                g = op["gradient"]
                cols = g.get("colors") if isinstance(g, dict) else None
                if not isinstance(cols, list) or len(cols) != 2 or not all(isinstance(c, str) for c in cols):
                    errors.append(
                        '%s: gradient needs {"colors": ["RRGGBB", "RRGGBB"]} (exactly 2), optional "angle" (degrees) and "positions" ([0-1, 0-1])'
                        % tag
                    )
    return errors


def collect_issue_map(path, only_slides=None):
    """(slide_idx, sid) -> issues dict, measured on a freshly-loaded instance."""
    prs = Presentation(path)
    index = build_index(prs, measure=True, only_slides=only_slides)
    out = {}
    for slide_idx, recs in index.items():
        for sid, r in recs.items():
            iss = rec_issues(r)
            if iss:
                out[(slide_idx, sid)] = iss
    return out


def cmd_apply(args):
    patch = json.loads(Path(args.patch).read_text(encoding="utf-8"))
    ops = patch["ops"] if isinstance(patch, dict) else patch
    if not isinstance(ops, list) or not ops:
        sys.exit("Error: patch must be {\"ops\": [...]} or a bare op array (got empty/invalid)")
    out_path = resolve_output(args)

    before = collect_issue_map(args.file)

    prs = Presentation(args.file)
    ctx = Ctx(prs, build_index(prs, measure=False))
    errors = validate_ops(ctx, ops)
    if errors:
        print("PATCH REJECTED — %d validation error(s), nothing was modified:" % len(errors))
        for e in errors:
            print("  - " + e)
        sys.exit(1)

    for i, op in enumerate(ops):
        if "slide" in op:
            ctx.touched.add(op["slide"])
        try:
            OP_HANDLERS[op["op"]](ctx, op)
        except PatchError as e:
            print("PATCH FAILED at op[%d] (%s): %s" % (i, op["op"], e))
            print("Nothing was saved.")
            sys.exit(1)

    prs.save(out_path)

    after = collect_issue_map(out_path)
    new_issues = []
    for key, iss in after.items():
        prev = before.get(key, {})
        for k, v in iss.items():
            if k == "covered_by":
                if k not in prev:
                    new_issues.append(
                        "slide %d %s: text extends under picture(s) %s — it renders CLIPPED behind them"
                        % (key[0], key[1], v))
            elif k in ("overlaps", "warnings"):
                if k not in prev:
                    new_issues.append("slide %d %s: new %s %s" % (key[0], key[1], k, v))
            elif isinstance(v, (int, float)) and v > prev.get(k, 0) + 0.05:
                new_issues.append('slide %d %s: %s %.2f" (was %.2f")' % (key[0], key[1], k, v, prev.get(k, 0)))

    report = {
        "saved": str(out_path),
        "ops_applied": len(ops),
        "log": ctx.log,
        "touched_slides": sorted(ctx.touched),
        "new_issues": new_issues,
    }
    if args.json:
        print(json.dumps(report, indent=2, ensure_ascii=False))
    else:
        print("Applied %d op(s) -> %s" % (len(ops), out_path))
        for line in ctx.log:
            print("  ✓ " + line)
        if new_issues:
            print("⚠ new/worsened geometry issues (run `deck.py %s fix --slides %s`):" % (out_path, ",".join(str(s) for s in sorted(ctx.touched))))
            for line in new_issues:
                print("  - " + line)

    # Optional powerups — chain the rest of the edit loop in this same call.
    # Both remain fully available as standalone subcommands.
    touched_csv = ",".join(str(s) for s in sorted(ctx.touched))
    if getattr(args, "fix", False):
        if touched_csv:
            print("\n--fix: deterministic repair on touched slides (%s)" % touched_csv)
            cmd_fix(argparse.Namespace(
                file=str(out_path), slides=touched_csv, all=False,
                output=None, in_place=True, json=False,
            ))
        else:
            print("\n--fix: no slide-level edits to fix (master-only patch?) — skipped")
    if getattr(args, "render", None):
        if touched_csv:
            print("\n--render: refreshing touched slides -> %s" % args.render)
            cmd_render(argparse.Namespace(
                file=str(out_path), output=args.render, slide=touched_csv,
                dpi=110, crop=None, scale=1,
            ))
        else:
            print("\n--render: no touched slides — skipped")


# ---------------------------------------------------------------------------
# fix — deterministic geometry repair
# ---------------------------------------------------------------------------

MIN_FONT_SCALE = 0.6
MIN_FONT_PT = 10.0
SLIDE_MARGIN_IN = 0.1
BG_AREA_FRACTION = 0.6  # shapes covering >60% of the slide don't block growth


def _effective_font_pt(para, fallback):
    pd = ParagraphData(para)
    if pd.font_size:
        return pd.font_size
    return fallback


def _scale_font(shape, scale, default_pt):
    changed = []
    for para in shape.text_frame.paragraphs:
        if not para.text.strip():
            continue
        eff = _effective_font_pt(para, default_pt)
        new = max(MIN_FONT_PT, round(eff * scale * 2) / 2)
        for run in para.runs:
            run.font.size = Pt(new)
        changed.append((eff, new))
    return changed


def cmd_fix(args):
    if not args.slides and not args.all:
        sys.exit(
            "Error: fix needs an explicit scope: --slides 3,7 (the slides you edited) "
            "or --all. This avoids 'fixing' measurement false-positives on slides you never touched."
        )
    out_path = resolve_output(args)
    only = None if args.all else parse_slide_list(args.slides)

    # measure on a throwaway instance, edit a clean one
    measured = Presentation(args.file)
    midx = build_index(measured, measure=True, only_slides=only)
    prs = Presentation(args.file)
    eidx = build_index(prs, measure=False, only_slides=only)

    slide_w = inches(prs.slide_width)
    slide_h = inches(prs.slide_height)
    slide_area = slide_w * slide_h

    fixed, residue = [], []
    for slide_idx, recs in midx.items():
        for sid, r in recs.items():
            iss = rec_issues(r)
            if not iss:
                continue
            er = eidx[slide_idx].get(sid)
            if er is None:
                continue
            sh = er.shape

            ov = iss.get("frame_overflow_bottom")
            if ov is not None and ov > 0.1:
                # (a) grow if free space below
                needed = ov + 0.05
                bottom = r.top + r.height
                blocked = False
                for osid, o in midx[slide_idx].items():
                    if osid == sid or osid == r.group or o.group == sid:
                        continue
                    if o.width * o.height > BG_AREA_FRACTION * slide_area:
                        continue  # background art doesn't block
                    h_overlap = min(r.left + r.width, o.left + o.width) - max(r.left, o.left)
                    v_intersect = (o.top < bottom + needed + 0.1) and (o.top + o.height > bottom)
                    if h_overlap > 0.1 and v_intersect:
                        blocked = True
                        break
                if not blocked and bottom + needed <= slide_h - SLIDE_MARGIN_IN:
                    sh.height = (sh.height or 0) + int(Inches(needed) / er.tf.sy)
                    fixed.append({"slide": slide_idx, "shape": sid, "action": "grow", "was": ov,
                                  "detail": 'height +%.2f" (text overflowed %.2f")' % (needed, ov)})
                else:
                    # (b) shrink-to-fit via explicit font scaling
                    usable = max(r.height - 0.1, 0.2)
                    scale = max(usable / (usable + ov), MIN_FONT_SCALE)
                    default_pt = (r.sd.default_font_size or r.sd._get_default_font_size()) if r.sd else 14
                    changed = _scale_font(sh, scale, default_pt)
                    if not any(abs(a - b) > 0.25 for a, b in changed):
                        # floors made it a no-op — don't pretend we fixed anything
                        residue.append({
                            "slide": slide_idx, "shape": sid,
                            "issue": 'text overflows %.2f" but the font is already at the readability floor (%gpt)'
                                     % (ov, MIN_FONT_PT),
                            "suggest": "shorten the text with set-text, or resize/move the shape",
                        })
                    else:
                        fixed.append({"slide": slide_idx, "shape": sid, "action": "shrink-font", "was": ov,
                                      "detail": "font x%.2f (%s)" % (scale,
                                          ", ".join("%g->%g" % (a, b) for a, b in changed[:4])
                                          + ("…" if len(changed) > 4 else ""))})

            for key, axis in (("slide_overflow_right", "x"), ("slide_overflow_bottom", "y")):
                ovs = iss.get(key)
                if ovs is None:
                    continue
                if not r.is_text:
                    # pictures off one edge MAY be intentional bleed — never auto-move
                    residue.append({
                        "slide": slide_idx, "shape": sid,
                        "issue": '%s (%s) extends %.2f" off the slide (may be intentional bleed)' % (sid, r.type, ovs),
                        "suggest": "move %s by [%s] if unintended" % (sid, ("-%.2f, 0" % ovs) if axis == "x" else ("0, -%.2f" % ovs)),
                    })
                    continue
                if axis == "x":
                    new_left = max(0.0, r.left - ovs)
                    sh.left = int(er.tf.child_x(Inches(new_left)))
                    if new_left == 0.0 and r.width > slide_w:
                        sh.width = int(Inches(slide_w) / er.tf.sx)
                        fixed.append({"slide": slide_idx, "shape": sid, "action": "fit-width", "was": ovs,
                                      "detail": "moved to x=0 and shrunk to slide width"})
                    else:
                        fixed.append({"slide": slide_idx, "shape": sid, "action": "nudge-left", "was": ovs,
                                      "detail": '%.2f" back onto the slide' % ovs})
                else:
                    new_top = max(0.0, r.top - ovs)
                    sh.top = int(er.tf.child_y(Inches(new_top)))
                    fixed.append({"slide": slide_idx, "shape": sid, "action": "nudge-up", "was": ovs,
                                  "detail": '%.2f" back onto the slide' % ovs})

            if "overlaps" in iss:
                for other_sid, area in iss["overlaps"].items():
                    if sid < other_sid:  # report each pair once
                        o = midx[slide_idx].get(other_sid)
                        suggest = ""
                        if o is not None:
                            dy = (r.top + r.height) - o.top
                            suggest = 'move %s by [0, %.2f]' % (other_sid, round(dy + 0.05, 2))
                        residue.append({
                            "slide": slide_idx, "shape": sid,
                            "issue": "overlaps %s by %.2f sq in (needs judgment — could be intentional design)" % (other_sid, area),
                            "suggest": suggest,
                        })

            if "covered_by" in iss:
                for pic_sid, area in iss["covered_by"].items():
                    residue.append({
                        "slide": slide_idx, "shape": sid,
                        "issue": "text extends under PICTURE %s by %.2f sq in — it renders clipped/hidden behind the picture" % (pic_sid, area),
                        "suggest": "reorder %s z:front to draw the text on top, or move/shrink them apart" % sid,
                    })

    if fixed:
        prs.save(out_path)
    elif str(out_path) != str(args.file):
        shutil.copy2(args.file, out_path)
    remaining = collect_issue_map(out_path, only_slides=only)

    # honesty pass: re-measure — anything still (nearly) as broken as before is residue, not fixed
    confirmed = []
    for f in fixed:
        rem = remaining.get((f["slide"], f["shape"]), {})
        key = "frame_overflow_bottom" if f["action"] in ("grow", "shrink-font") else (
            "slide_overflow_right" if f["action"] in ("nudge-left", "fit-width") else "slide_overflow_bottom")
        still = rem.get(key)
        if still is not None and still > 0.12:
            residue.append({
                "slide": f["slide"], "shape": f["shape"],
                "issue": '%s applied (%s) but %s is still %.2f" (was %.2f")' % (f["action"], f["detail"], key, still, f.get("was", 0)),
                "suggest": "shorten the text with set-text, or resize/move the shape",
            })
        else:
            confirmed.append(f)
    fixed = confirmed

    report = {"saved": str(out_path), "fixed": fixed, "residue": residue,
              "remaining_issue_shapes": ["slide %d %s: %s" % (k[0], k[1], v) for k, v in remaining.items()]}
    if args.json:
        print(json.dumps(report, indent=2, ensure_ascii=False))
    else:
        print("fix -> %s: %d fixed, %d residue" % (out_path, len(fixed), len(residue)))
        for f in fixed:
            print("  ✓ slide %(slide)d %(shape)s: %(action)s — %(detail)s" % f)
        for rsd in residue:
            print("  ⚠ slide %(slide)d %(shape)s: %(issue)s" % rsd + ("  → suggest: %s" % rsd["suggest"] if rsd.get("suggest") else ""))
        if remaining:
            print("  remaining measured issues:")
            for k, v in remaining.items():
                print("    slide %d %s: %s" % (k[0], k[1], v))


# ---------------------------------------------------------------------------
# render
# ---------------------------------------------------------------------------

def _unhide_all(pptx_path):
    prs = Presentation(pptx_path)
    changed = False
    for s in prs.slides:
        if s.element.get("show") is not None:
            del s.element.attrib["show"]
            changed = True
    if changed:
        prs.save(pptx_path)


def cmd_render(args):
    outdir = Path(args.output)
    outdir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(args.file)
    n = len(prs.slides._sldIdLst)
    hidden = [i for i, s in enumerate(prs.slides) if s.element.get("show") == "0"]

    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        if args.slide:
            wanted = sorted(parse_slide_list(args.slide))
            bad = [i for i in wanted if not (0 <= i < n)]
            if bad:
                sys.exit("Error: slide index(es) %s out of range (0-%d) — 0-BASED" % (bad, n - 1))
            sub = td / "subset.pptx"
            from rearrange import rearrange_presentation

            with contextlib.redirect_stdout(io.StringIO()):
                rearrange_presentation(Path(args.file), sub, wanted)
            _unhide_all(sub)  # subset render shows even hidden slides
            target, mapping = sub, wanted
        else:
            target = td / "full.pptx"
            shutil.copy2(args.file, target)
            mapping = [i for i in range(n) if i not in hidden]
            if hidden:
                print("note: hidden slide(s) %s are skipped in a full render (render them explicitly with --slide)" % hidden)

        r = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", str(target), "--outdir", str(td)],
            capture_output=True, text=True,
        )
        pdf = target.with_suffix(".pdf")
        if not pdf.exists():
            sys.exit("Error: soffice PDF conversion failed:\n%s%s" % (r.stdout, r.stderr))
        r = subprocess.run(
            ["pdftoppm", "-jpeg", "-r", str(args.dpi), str(pdf), str(td / "page")],
            capture_output=True, text=True,
        )
        pages = sorted(td.glob("page-*.jpg"), key=lambda p: int(p.stem.split("-")[-1]))
        if not pages:
            sys.exit("Error: pdftoppm produced no images:\n%s%s" % (r.stdout, r.stderr))
        if len(pages) != len(mapping):
            print(
                "warning: %d rendered pages vs %d expected slides — mapping the first %d"
                % (len(pages), len(mapping), min(len(pages), len(mapping)))
            )
        written = []
        for page, idx in zip(pages, mapping):
            dest = outdir / ("slide-%d.jpg" % idx)
            shutil.move(str(page), dest)
            written.append((dest, idx))
    if args.crop:
        # crop is in INCHES so coordinates line up with inspect geometry
        try:
            l, t, w, h = (float(x) for x in args.crop.split(","))
        except ValueError:
            sys.exit("Error: --crop wants l,t,w,h in inches, e.g. --crop 0.3,2.0,6.5,1.5")
        from PIL import Image

        slide_w_in = inches(prs.slide_width)
        for dest, idx in list(written):
            with Image.open(dest) as im:
                ppi = im.width / slide_w_in
                box = (
                    max(0, int(l * ppi)),
                    max(0, int(t * ppi)),
                    min(im.width, int((l + w) * ppi)),
                    min(im.height, int((t + h) * ppi)),
                )
                if box[2] <= box[0] or box[3] <= box[1]:
                    sys.exit("Error: --crop region %s lies outside the slide" % args.crop)
                crop = im.crop(box)
                if args.scale != 1:
                    crop = crop.resize(
                        (int(crop.width * args.scale), int(crop.height * args.scale)),
                        Image.LANCZOS,
                    )
                cdest = outdir / ("slide-%d-crop.jpg" % idx)
                crop.save(cdest, quality=90)
                written.append((cdest, idx))
    print(
        "Rendered %d image(s) -> %s: %s"
        % (len(written), outdir, ", ".join(sorted(d.name for d, _ in written)))
    )


# ---------------------------------------------------------------------------
# diff
# ---------------------------------------------------------------------------

def _media_crcs(path):
    with zipfile.ZipFile(path) as z:
        return {n.split("/")[-1]: z.getinfo(n).CRC for n in z.namelist() if n.startswith("ppt/media/")}


def _slide_text_map(recs):
    out = {}
    for sid, r in recs.items():
        if r.is_text:
            out[sid] = [p.text.strip() for p in r.shape.text_frame.paragraphs if p.text.strip()]
        elif r.is_table:
            out[sid] = [c.text for row in r.shape.table.rows for c in row.cells]
    return out


def cmd_diff(args):
    if not Path(args.other).exists():
        sys.exit("Error: file not found: %s" % args.other)
    pa, pb = Presentation(args.file), Presentation(args.other)
    na, nb = len(pa.slides._sldIdLst), len(pb.slides._sldIdLst)
    lines = []
    if na != nb:
        lines.append("slide count: %d -> %d" % (na, nb))
    crc_a, crc_b = _media_crcs(args.file), _media_crcs(args.other)
    ia = build_index(pa, measure=False)
    ib = build_index(pb, measure=False)
    for i in range(min(na, nb)):
        slide_lines = []
        ra, rb = ia.get(i, {}), ib.get(i, {})
        ta, tb = _slide_text_map(ra), _slide_text_map(rb)
        for sid in sorted(set(ta) | set(tb), key=lambda s: (len(s), s)):
            if sid not in tb:
                slide_lines.append("  - %s removed (text was: %s)" % (sid, " / ".join(ta[sid])[:60]))
            elif sid not in ta:
                slide_lines.append("  + %s added: %s" % (sid, " / ".join(tb[sid])[:60]))
            elif ta[sid] != tb[sid]:
                slide_lines.append("  ~ %s text: %s -> %s" % (sid, " / ".join(ta[sid])[:50], " / ".join(tb[sid])[:50]))
        for sid in sorted(set(ra) & set(rb), key=lambda s: (len(s), s)):
            a, b = ra[sid], rb[sid]
            if abs(a.left - b.left) > 0.03 or abs(a.top - b.top) > 0.03:
                slide_lines.append("  ~ %s moved [%.2f,%.2f] -> [%.2f,%.2f]" % (sid, a.left, a.top, b.left, b.top))
            if abs(a.width - b.width) > 0.03 or abs(a.height - b.height) > 0.03:
                slide_lines.append("  ~ %s resized %sx%s -> %sx%s" % (sid, a.width, a.height, b.width, b.height))
            ma = {rid: m for rid, m in a.rids}
            mb = {rid: m for rid, m in b.rids}
            if set(ma.values()) != set(mb.values()):
                slide_lines.append("  ~ %s image %s -> %s" % (sid, sorted(ma.values()), sorted(mb.values())))
            else:
                for m in set(ma.values()):
                    if crc_a.get(m) != crc_b.get(m):
                        slide_lines.append("  ~ %s media %s bytes changed" % (sid, m))
        for sid in sorted(set(rb) - set(ra), key=lambda s: (len(s), s)):
            if sid not in tb:  # non-text additions (text ones already reported)
                slide_lines.append("  + %s added (%s)" % (sid, rb[sid].type))
        for sid in sorted(set(ra) - set(rb), key=lambda s: (len(s), s)):
            if sid not in ta:
                slide_lines.append("  - %s removed (%s)" % (sid, ra[sid].type))
        # notes
        try:
            na_t = pa.slides[i].notes_slide.notes_text_frame.text if pa.slides[i].has_notes_slide else ""
            nb_t = pb.slides[i].notes_slide.notes_text_frame.text if pb.slides[i].has_notes_slide else ""
            if na_t.strip() != nb_t.strip():
                slide_lines.append("  ~ notes changed (%d -> %d chars)" % (len(na_t.strip()), len(nb_t.strip())))
        except Exception:
            pass
        # transition
        tra, trb = _transition_dict(pa.slides[i]), _transition_dict(pb.slides[i])
        if tra != trb:
            fmt = lambda t: "none" if not t else " ".join("%s=%s" % (k, t[k]) for k in sorted(t))
            slide_lines.append("  ~ transition %s -> %s" % (fmt(tra), fmt(trb)))
        if slide_lines:
            lines.append("slide %d:" % i)
            lines.extend(slide_lines)
    if not lines:
        print("No structural differences.")
    else:
        print("\n".join(lines))


# ---------------------------------------------------------------------------
# slides / merge — deck-level structure (wrap rearrange.py / merge_decks.py)
# ---------------------------------------------------------------------------

def cmd_slides(args):
    try:
        seq = [int(x.strip()) for x in args.sequence.split(",")]
    except ValueError:
        sys.exit("Error: sequence must be comma-separated 0-based indices, e.g. 0,3,3,5")
    out_path = resolve_output(args)
    from rearrange import rearrange_presentation

    with contextlib.redirect_stdout(io.StringIO()):
        try:
            rearrange_presentation(Path(args.file), Path(out_path), seq)
        except ValueError as e:
            sys.exit("Error: %s" % e)
    print("slides -> %s: kept [%s] in that order (%d slides; repeats were duplicated)"
          % (out_path, args.sequence, len(seq)))


def cmd_merge(args):
    from merge_decks import list_layouts, merge

    if args.list_layouts:
        list_layouts(Presentation(args.file))
        return
    if not args.module:
        sys.exit("Error: merge needs the source .pptx: deck.py base.pptx merge module.pptx -o out.pptx")
    if not Path(args.module).exists():
        sys.exit("Error: file not found: %s" % args.module)
    out_path = resolve_output(args)
    slide_nums = None
    if args.slides:
        # deck.py speaks 0-based everywhere; merge_decks.merge wants 1-based
        slide_nums = [int(x.strip()) + 1 for x in args.slides.split(",")]
    buf = io.StringIO()
    with contextlib.redirect_stdout(buf):
        merge(Path(args.file), Path(out_path), Path(args.module), slide_nums, args.at, args.layout)
    n_in = len(slide_nums) if slide_nums else len(Presentation(args.module).slides._sldIdLst)
    n_out = len(Presentation(out_path).slides._sldIdLst)
    print("merge -> %s: %d slide(s) from %s inserted at %s (%d slides total)"
          % (out_path, n_in, Path(args.module).name,
             ("end" if args.at is None else args.at), n_out))
    for line in buf.getvalue().splitlines():
        if "⚠" in line or "warning" in line.lower():
            print("  " + line.strip())


# ---------------------------------------------------------------------------
# xml — the escape hatch, tooled
# ---------------------------------------------------------------------------

def _part_name(args):
    if args.part:
        return args.part
    if args.slide is not None:
        return "ppt/slides/slide%d.xml" % (args.slide + 1)
    sys.exit("Error: xml needs --slide N (0-based) or --part ppt/...xml")


def cmd_xml(args):
    from lxml import etree

    part = _part_name(args)
    if args.action == "get":
        with zipfile.ZipFile(args.file) as z:
            if part not in z.namelist():
                near = [n for n in z.namelist() if n.startswith("ppt/")][:40]
                sys.exit("Error: %s not in package. Parts include:\n  %s" % (part, "\n  ".join(near)))
            data = z.read(part)
        tree = etree.fromstring(data)
        pretty = etree.tostring(tree, pretty_print=True, xml_declaration=True, encoding="UTF-8")
        out = Path(args.output) if args.output else Path(Path(part).name)
        out.write_bytes(pretty)
        print("Wrote pretty-printed %s -> %s" % (part, out))
        print("Edit it, then write back with: deck.py %s xml set %s %s -o out.pptx"
              % (args.file, out, ("--slide %d" % args.slide) if args.slide is not None else "--part " + part))
    else:  # set
        if not args.xmlfile:
            sys.exit("Error: xml set needs the edited XML file right after 'set': xml set <file.xml> --slide N")
        if not Path(args.xmlfile).exists():
            sys.exit("Error: XML file not found: %s" % args.xmlfile)
        raw = Path(args.xmlfile).read_bytes()
        try:
            etree.fromstring(raw)
        except etree.XMLSyntaxError as e:
            sys.exit("Error: edited XML does not parse: %s" % e)
        out_path = resolve_output(args)
        tmp = str(out_path) + ".tmp"
        with zipfile.ZipFile(args.file) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            if part not in zin.namelist():
                sys.exit("Error: %s not in package" % part)
            for item in zin.infolist():
                data = raw if item.filename == part else zin.read(item.filename)
                zout.writestr(item, data)
        shutil.move(tmp, out_path)
        try:
            Presentation(out_path)  # sanity: python-pptx must reopen it
        except Exception as e:
            sys.exit("Error: package no longer opens after edit (%s). Output left at %s for inspection." % (e, out_path))
        print("Wrote %s into %s -> %s (package reopens OK)" % (args.xmlfile, part, out_path))
        if args.slide is not None:
            issues = collect_issue_map(out_path, only_slides={args.slide})
            if issues:
                for k, v in issues.items():
                    print("  ⚠ slide %d %s: %s" % (k[0], k[1], v))


# ---------------------------------------------------------------------------
# docs — the full reference, so agents never have to read this file's source
# ---------------------------------------------------------------------------

DOCS = """\
deck.py REFERENCE (print anytime: `deck.py docs` — no file needed)

GOLDEN RULES
- ALL slide indices are 0-BASED everywhere (slide 0 = first). Rendered images
  are named slide-<0-based>.jpg to match.
- Shapes are addressed by NATIVE id: "s12". Ids are stable — deleting a shape
  never renumbers the others. Get ids from `inspect`. Exact shape NAME also
  works if unique on the slide.
- All positions/sizes are INCHES, measured from the slide's top-left.
- Patches are ATOMIC: every op is validated first (all errors reported at
  once, listing the slide's real shapes), then applied all-or-nothing.
- After apply, the tool lints and reports only NEW/WORSENED geometry issues,
  with the exact `fix` command to run. A `covered_by` report (text under a
  picture) means the text RENDERS CLIPPED — treat it as real, not stylistic.

THE EDIT LOOP
  inspect --slide N  →  write patch  →  apply  →  fix --slides <touched>
  →  render --slide <touched> (look at it!)  →  diff (optional sanity)

  ONE-CALL POWERUP: `apply patch.json --in-place --fix --render img/` chains
  apply + fix + render-of-touched-slides in a single call. fix/render remain
  standalone subcommands — use them individually whenever you want.

  ORIENTATION READS: `inspect --slide N --brief` = one line per shape (id,
  type, geometry, font, text preview, issues). Use full JSON only when you
  are about to write a patch and need exact paragraph formatting.

PATCH FORMAT  {"ops":[ ... ]}  — ops run in order.

set-text       {"op":"set-text","slide":3,"shape":"s12","text":["Line one","Line two"]}
  - INHERITANCE: new paragraph i inherits ALL formatting of old paragraph i
    (font, size, bold, color, bullet, level, alignment, spacing). If you pass
    MORE paragraphs than existed, extras inherit from the LAST old paragraph
    (so adding bullets to a list continues the list style).
  - To CHANGE formatting, pass an object; your keys override the inherited:
      {"text":"Heading","bold":true,"font_size":24,"font_name":"Inter",
       "italic":false,"underline":false,"color":"D4A853","theme_color":"DARK_1",
       "alignment":"CENTER",            // LEFT|CENTER|RIGHT|JUSTIFY
       "bullet":true,"level":0,         // bullets ARE supported — never put
                                        // "•" in the text; level needs bullet
       "space_before":10,"space_after":6,"line_spacing":22}
  - To REMOVE a bullet inherited from the old paragraph: {"text":"...","bullet":false}.
  - MIXED FORMATTING inside one paragraph: pass "runs" instead of "text":
      {"runs":[{"text":"We measured "},
               {"text":"5x","bold":true,"color":"F8DE6E"},
               {"text":" fewer tokens"}]}
    Every run inherits the OLD paragraph's first-run font, then its own keys
    override (per-run keys: bold, italic, underline, font_size, font_name,
    color, theme_color, link). "link":"https://…" makes the run a hyperlink
    (null/"" removes one); works on plain paragraph objects too. Paragraph
    keys (alignment, bullet, level, spacing) still go on the paragraph
    object, next to "runs".
    A plain "text" paragraph still collapses mixed old runs to the first
    run's format — use "runs" when you need to PRESERVE a mid-paragraph style
    (inspect shows the current runs breakdown when formats differ).
  - TABLES: target a cell with "cell":[row,col] (0-based).
  - Keep text length comparable to the original or run `fix` after.

swap-image     {"op":"swap-image","slide":4,"rid":"rId3","image":"/abs/new.png"}
  - Target by "rid" (from inspect) or "shape":"s9". New image becomes its own
    media part: other slides sharing the old media are NOT affected.
  - "fit":"auto" (default) re-fits the frame, preserving aspect ratio, when it
    differs >8% from the frame; "stretch" fills the frame; "contain" forces fit.
  - A picture with a circular/rounded crop keeps that crop after the swap.
  - GLOBAL: {"op":"swap-image","media":"image13.png","image":"/abs/new.png"}
    (no slide/rid) overwrites that media file's bytes, so EVERY slide, master,
    and layout showing it changes at once — THE way to swap a logo deck-wide.
    Frames keep their geometry: supply art with the same aspect ratio.
  - Cannot fix client data baked INTO pixels — that needs a new asset.

replace-text   {"op":"replace-text","scope":"deck","from":"Globex","to":"Acme"}
  - scope "deck" = every slide; "master" = slide masters + layouts (footers
    live there); "slide" + "slide":N = one slide.
  - Replaces within single runs. If the text is split across runs, the error
    names the affected slides — use set-text on those shapes instead.
  - Errors if zero occurrences (so a silent no-op can't pass review).

replace-color  {"op":"replace-color","scope":"deck","from":"E8A33D","to":"F8DE6E"}
  - The re-theme primitive: swaps one concrete color for another everywhere
    in scope — fills, gradients, lines, text, effects. Same scopes as
    replace-text. One op per palette mapping; a full re-theme is a handful
    of ops in one atomic patch.
  - Only literal (srgbClr) colors match. Theme-indexed colors (schemeClr)
    are deliberately untouched — remap those once in the theme via xml.
  - Errors if zero occurrences, and the error lists the colors actually
    present (with counts) so the next patch can be exact.
  - Colors inside IMAGES don't change (pixels aren't XML) — re-render and
    look; swap-image any asset that carries the old palette.

set-props      {"op":"set-props","title":"Q3 Review","author":"Acme"}
  - Document metadata (File > Info): title, subject, author, keywords,
    comments, category, last_modified_by. Strings; "" clears a field.
    Current values appear in `inspect` (no --slide) under "props".

set-slide      {"op":"set-slide","slide":3,"hidden":true,"background":"0F5258"}
  - Slide-level properties. "hidden": true/false (hidden slides don't
    present/render unless explicitly listed). "background" paints the
    slide's own p:bg layer (under every shape): "RRGGBB",
    {"gradient":{"colors":["A","B"],"angle":90}}, or {"image":"/abs/p.png"}.
  - "transition": {"type":"fade"} sets how the slide ENTERS; "none" removes.
    Types: fade, cut, dissolve, push, wipe, split, cover, uncover, zoom.
    Options: "speed" slow|med|fast; "advance_after" seconds (auto-advance;
    click still advances unless "advance_on_click":false); per-type
    direction "dir" — push/wipe: l r u d; cover/uncover: + ld lu rd ru;
    split: "orient" horz|vert + "dir" in|out; zoom: "dir" in|out.
    One op per slide; same dict on every slide for a uniform deck.
    Transitions are MOTION — renders are static JPGs, so verify with
    inspect ("_transition" per slide) or diff, not render.

set-theme      {"op":"set-theme","colors":{"accent1":"BB7B19"},"fonts":{"major":"Georgia"}}
  - Remaps the THEME: color slots dk1 lt1 dk2 lt2 accent1-6 hlink folHlink,
    and major (headings) / minor (body) latin fonts. Every shape that uses
    theme-indexed colors/fonts follows automatically — this is how
    template decks are rebranded. "master":N targets one master (default
    all). Literal hard-coded colors don't follow — that's replace-color;
    a full rebrand is usually set-theme + a few replace-color ops.

set-notes      {"op":"set-notes","slide":3,"notes":"plain text speaker notes"}
move           {"op":"move","slide":3,"shape":"s12","to":[1.0,2.5]}     or "by":[dx,dy]
resize         {"op":"resize","slide":3,"shape":"s12","size":[4.0,1.5]} or "scale":0.8
set-style      {"op":"set-style","slide":3,"shape":"s12",...}
  - Text: applies to ALL runs in the shape (and all table cells): font_size,
    font_name, bold, italic, underline, color ("RRGGBB"). For per-paragraph
    changes use set-text override objects; per-run, set-text "runs".
  - Fill: "fill":"RRGGBB" = solid background; "fill":"none" = no fill
    (hollow shape — only the outline paints).
    "gradient":{"colors":["0F5258","4999A0"],"angle":90,"positions":[0,1]}
    = two-stop linear gradient (replaces any existing fill; angle in degrees
    counter-clockwise from left-to-right; positions optional, 0-1).
  - Line/border: "line_color":"RRGGBB", "line_width":1.5 (pt),
    "line_dash":"dash" (solid|dash|dot|dash_dot|dash_dot_dot|long_dash|
    long_dash_dot|round_dot|square_dot), "line":"none" removes the outline.
    Works on any shape incl. pictures (frame border).
  - "rotation":45 = degrees clockwise (inspect reports it; pos/size stay the
    unrotated bounding box).
  - "insets":[l,t,r,b] = text-frame internal margins in INCHES (PowerPoint
    defaults are 0.1/0.05/0.1/0.05, NOT zero — set [0,0,0,0] for flush text).
  - "adjustments":[0.12] = shape adjustment handles (e.g. roundRect corner
    radius as a fraction of the smaller dimension, max 0.5).
  - "shadow":false = remove the theme's inherited drop shadow (new autoshapes
    often carry one); true restores inheritance.
  - "alt_text":"describe the image" = accessibility description (any shape;
    "" removes). inspect shows it.
delete         {"op":"delete","slide":3,"shape":"s12"}   (deleting a GROUP deletes its children)
duplicate      {"op":"duplicate","slide":3,"shape":"s12","offset":[0,1.2],"text":["New label"]}
  - THE way to scale a styled element ("three boxes → four"). The copy keeps
    all styling/images, gets fresh ids, lands offset inches from the original;
    "text" then rewrites it via set-text semantics. "at":[l,t] = absolute.
copy-shape     {"op":"copy-shape","from_slide":8,"shape":"s12","slide":3,"at":[1.0,2.0],"text":["…"]}
  - Borrow a styled shape from ANOTHER slide; images/hyperlinks re-homed.
reorder        {"op":"reorder","slide":3,"shape":"s12","z":"front"}   front|back|forward|backward
  - Z-order. inspect lists shapes back-to-front (first = bottom layer).

CREATE OPS (prefer duplicate/copy-shape when a styled donor exists — new
shapes start from PowerPoint defaults, not the deck's design language.
Designing a whole slide from scratch? html2patch.py compiles an HTML/CSS
file into a patch of these ops: write the slide as HTML, get measured
add-shape/add-picture/add-table ops back. Needs playwright.)
add-shape      {"op":"add-shape","slide":3,"kind":"textbox","at":[1.0,2.0],"size":[4.0,1.5],
                "text":["Label"],"fill":"0B3D3A","font_size":18,"rotation":0,"name":"my-box"}
  - kind: textbox | rect | rounded_rect | ellipse | any MSO_SHAPE name
    (CHEVRON, PENTAGON, …) | line (straight connector, takes "from":[x,y] +
    "to":[x,y] instead of at/size).
  - Accepts every set-style key in the same op (fill, gradient, line_color,
    line_width, line_dash, font_*, color, rotation, insets, adjustments) plus
    "text" (set-text semantics — strings or {"text":…, formatting} objects or
    "runs").
  - Give it a "name" and LATER ops in the same patch can target that name
    (validation knows it's coming). New shape ids appear in the apply log.
add-picture    {"op":"add-picture","slide":3,"image":"/abs/img.png","at":[1.0,2.0],"width":4.0}
  - Give "width" OR "height" to keep aspect (preferred), "size":[w,h] to force
    both, neither = native size.
  - "crop":[l,t,r,b] = fractions of the source trimmed from each edge (this is
    how you fill a box without distortion: size to the box + crop the excess).
add-table      {"op":"add-table","slide":3,"at":[0.5,1.5],"size":[9.0,3.0],
                "rows":[["Header A","Header B"],["1","2"]],"font_size":14}
  - Optional: "name" (target it later in the same patch), "color":"RRGGBB"
    (all text), "fill":"RRGGBB"|"none" (all cells) or "fills":[[...]] row-major
    per-cell grid, "col_widths":[4.5,2.25,2.25] (inches), "first_row":false +
    "banding":false to neutralize the theme's banded table style.
  - set-text with "cell":[row,col] styles individual cells afterwards; its
    paragraph objects take "bullet":true (•) or "bullet":"number" (1. 2. 3.).
add-slide      {"op":"add-slide","layout":"Blank","at":5}
  - layout = name (exact, then substring) or 0-based index across all masters;
    omit it for the blankest layout. Omit "at" to append. Later ops in the
    same patch use POST-insertion indices; prefer adding slides at the end or
    in a separate patch to keep indices easy to reason about.

TABLE STRUCTURE (tables without merged cells only — merged → xml escape hatch)
add-row        {"op":"add-row","slide":3,"shape":"s12","cells":["a","b"],"copy":0,"at":2}
  - Clones row "copy" (default: last) with all its formatting, inserts at
    "at" (default: append), then writes "cells" via set-text inheritance.
delete-row     {"op":"delete-row","slide":3,"shape":"s12","row":2}      (negative = from end)
add-col        {"op":"add-col","slide":3,"shape":"s12","cells":["h","1"],"copy":-1,"at":1}
  - Same semantics column-wise; all column widths rescale so the table keeps
    its overall width. "cells" runs top-to-bottom (one per row).
delete-col     {"op":"delete-col","slide":3,"shape":"s12","col":1}      (widths rescale back)

FIX (deterministic repair — run after every apply)
  fix --slides 3,7 --in-place        # ALWAYS scope to slides you touched
  - Text overflow: grows the box if free space below, else shrinks fonts
    (floor: 60% of original and 10pt). Off-slide TEXT: nudged back.
  - NEVER auto-moves pictures (off-slide may be intentional bleed) and never
    auto-resolves overlaps — those come back as residue with a suggested op.
  - Residue = what still needs judgment. It is honest: an action that didn't
    actually resolve the measured issue is reported as residue, not success.
  - Why scoping: overflow is ESTIMATED with fallback fonts when the deck's
    fonts aren't installed; untouched template slides have known false
    positives. Never `--all` a template-derived deck.

INSPECT FIELDS
  pos/size [in], type, group (parent group id), placeholder,
  paragraphs (text + non-default formatting — what set-text will inherit;
    a "runs" breakdown appears when a paragraph mixes run formats;
    "link" appears on hyperlinked runs),
  rid + media (pictures), rows (tables), fill / fill_gradient / line (when set),
  alt_text (when set), "_fonts" + "_notes" + "_transition" per slide,
  deck-level: "props" (document metadata), "hidden_slides", issues:
    frame_overflow_bottom  text taller than its box (inches over)
    slide_overflow_right/bottom  shape sticks off the slide
    overlaps  {other_sid: sq inches}  (text-vs-text only)
    covered_by  {picture_sid: sq inches}  text drawn UNDER a picture that is
      above it in z-order — including its estimated overflow region. This
      text renders clipped/hidden; it is almost never a false positive.
  --issues = only problem shapes; --master = masters/layouts too (footers!)

DECK STRUCTURE (subcommands, not patch ops)
  slides 0,3,3,5 -o out.pptx   keep only these slides, in this order, 0-based;
    a repeated index duplicates that slide (styling + images preserved).
  merge module.pptx --slides 0,2 --at 12 --layout 5 -o out.pptx
    copy slides from another deck; they re-attach to THIS deck's layout
    (--list-layouts to choose) so they inherit its master/footer/theme.
    Charts/OLE relationships are not copied (warned, handle manually).

RENDER / DIFF / XML
  render -o dir [--slide 3,7] [--dpi 110] [--crop l,t,w,h --scale 2]
    crop is in INCHES (matches inspect coords); writes slide-<N>-crop.jpg.
    Hidden slides are skipped in full renders (warned), rendered when explicit.
  diff other.pptx   text/geometry/media/notes changelog — verify without rendering
  xml get --slide 5 -o s5.xml   pretty-printed part XML (escape hatch);
  xml set s5.xml --slide 5 -o out.pptx   parse-checked + lint-checked write-back.
  Use xml ONLY when no op above expresses the change.

OUT OF SCOPE (use the xml escape hatch, or do it in PowerPoint)
  Creating/editing native charts, shape ANIMATIONS (slide transitions ARE
  covered — set-slide "transition"), embedded video/OLE, merged table cells.
  swap-image DOES handle replacing a chart rendered as a picture; add-table
  covers most "chart as table" needs.

RECIPES
  Rebrand a deck:      replace-text scope master (footer) + scope deck (mentions)
  Re-theme a deck:     inspect the palette (xml get + grep srgbClr, or render),
                       then ONE patch of replace-color ops — one per mapping,
                       deck + master scope; render and check images for
                       old-palette pixels baked into pictures
  Scale a 3-item list to 5: two duplicates with offsets + text, then fix
  Make text fit:       prefer set-text with shorter copy; else set-style
                       font_size; else resize; fix handles the remainder
  Match fonts across merged slides: set-style font_name per shape
  Borrow a styled card: copy-shape from the slide that has it, then set-text
"""

# ---------------------------------------------------------------------------
# main
# ---------------------------------------------------------------------------

def resolve_output(args):
    if getattr(args, "in_place", False):
        return args.file
    if getattr(args, "output", None):
        Path(args.output).parent.mkdir(parents=True, exist_ok=True)
        return args.output
    sys.exit("Error: pass -o <out.pptx> or --in-place")


def main():
    # `deck.py docs` needs no file — print the reference and exit
    if len(sys.argv) > 1 and sys.argv[1] == "docs":
        print(DOCS)
        return

    parser = argparse.ArgumentParser(
        description="Unified deck editor. ALL slide indices are 0-BASED. Full reference: deck.py docs",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="Patch format" + __doc__.split("Patch format", 1)[1],
    )
    parser.add_argument("file", help="the .pptx to operate on")
    sub = parser.add_subparsers(dest="cmd", required=True)

    sub.add_parser("docs", help="print the full patch reference + semantics")

    p = sub.add_parser("inspect", help="JSON index of every shape (text, images, tables, issues)")
    p.add_argument("--slide", help="only these slides, e.g. 3,7 (0-based)")
    p.add_argument("--issues", action="store_true", help="only shapes with geometric issues")
    p.add_argument("--brief", action="store_true", help="compact one-line-per-shape text (orientation reads)")
    p.add_argument("--master", action="store_true", help="include slide masters + layouts (footers live there)")
    p.add_argument("-o", "--output", help="write JSON here instead of stdout")

    p = sub.add_parser("apply", help="apply a JSON patch of ops (see module docstring)")
    p.add_argument("patch", help="patch JSON file: {\"ops\": [...]}")
    p.add_argument("-o", "--output")
    p.add_argument("--in-place", action="store_true")
    p.add_argument("--json", action="store_true", help="machine-readable report")
    p.add_argument("--fix", action="store_true", help="powerup: chain `fix` on touched slides after applying")
    p.add_argument("--render", metavar="DIR", help="powerup: chain `render` of touched slides into DIR")

    p = sub.add_parser("fix", help="deterministic geometry repair (overflow, off-slide)")
    p.add_argument("--slides", help="scope: the slides you edited, e.g. 3,7 (0-based)")
    p.add_argument("--all", action="store_true", help="whole deck (beware template false-positives)")
    p.add_argument("-o", "--output")
    p.add_argument("--in-place", action="store_true")
    p.add_argument("--json", action="store_true")

    p = sub.add_parser("render", help="render slides to outdir/slide-<idx>.jpg (0-based names)")
    p.add_argument("-o", "--output", required=True, help="output directory")
    p.add_argument("--slide", help="only these slides, e.g. 3,7 (0-based)")
    p.add_argument("--dpi", type=int, default=110)
    p.add_argument("--crop", help="zoom a region: l,t,w,h in INCHES (matches inspect coords)")
    p.add_argument("--scale", type=float, default=1, help="upscale factor for the crop (e.g. 2)")

    p = sub.add_parser("diff", help="structural changelog vs another deck")
    p.add_argument("other", help="the other .pptx")

    p = sub.add_parser("slides", help="reorder/duplicate/delete slides by sequence (0-based)")
    p.add_argument("sequence", help="e.g. 0,3,3,5 — keep only these, in this order (repeat = duplicate)")
    p.add_argument("-o", "--output")
    p.add_argument("--in-place", action="store_true")

    p = sub.add_parser("merge", help="copy slides from another deck into this one")
    p.add_argument("module", nargs="?", help="source .pptx to pull slides from")
    p.add_argument("--slides", help="source slide indices, 0-based, e.g. 0,2,5 (default: all)")
    p.add_argument("--at", type=int, help="0-based insert position in this deck (default: append)")
    p.add_argument("--layout", type=int, help="layout index of THIS deck to attach imported slides to")
    p.add_argument("--list-layouts", action="store_true", help="list this deck's layouts and exit")
    p.add_argument("-o", "--output")
    p.add_argument("--in-place", action="store_true")

    p = sub.add_parser("xml", help="escape hatch: get/set pretty-printed part XML")
    p.add_argument("action", choices=["get", "set"])
    p.add_argument("xmlfile", nargs="?", help="(set) the edited XML file — must come right after 'set'")
    p.add_argument("--slide", type=int, help="slide index (0-based)")
    p.add_argument("--part", help="explicit part name, e.g. ppt/slideMasters/slideMaster1.xml")
    p.add_argument("-o", "--output")
    p.add_argument("--in-place", action="store_true")

    args = parser.parse_args()
    if not Path(args.file).exists():
        sys.exit("Error: file not found: %s" % args.file)

    {
        "docs": lambda a: print(DOCS),
        "inspect": cmd_inspect,
        "apply": cmd_apply,
        "fix": cmd_fix,
        "render": cmd_render,
        "diff": cmd_diff,
        "slides": cmd_slides,
        "merge": cmd_merge,
        "xml": cmd_xml,
    }[args.cmd](args)


if __name__ == "__main__":
    main()
