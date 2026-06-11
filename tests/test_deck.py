"""End-to-end tests for deck.py — drive the CLI exactly as an agent would.

No binary fixtures: every test builds its deck with python-pptx, then talks to
deck.py only through its command line and JSON patches.
"""
import json
import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

SCRIPTS = Path(__file__).resolve().parents[1] / "skills" / "hands-on-deck" / "scripts"
DECK = SCRIPTS / "deck.py"


def run(*args):
    return subprocess.run(
        [sys.executable, str(DECK), *map(str, args)], capture_output=True, text=True
    )


def inspect(path, *args):
    r = run(path, "inspect", *args)
    assert r.returncode == 0, r.stderr
    return json.loads(r.stdout)


def apply_patch(deck_path, ops, out, *extra):
    patch = deck_path.parent / "patch.json"
    patch.write_text(json.dumps({"ops": ops}))
    return run(deck_path, "apply", patch, "-o", out, *extra)


@pytest.fixture
def img(tmp_path):
    p = tmp_path / "pic.png"
    Image.new("RGB", (200, 120), (40, 90, 100)).save(p)
    return p


@pytest.fixture
def deck(tmp_path, img):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
    blank = prs.slide_layouts[6]

    s0 = prs.slides.add_slide(blank)
    tb = s0.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    run0 = tb.text_frame.paragraphs[0].add_run()
    run0.text = "Hello World"
    run0.font.size, run0.font.bold = Pt(30), True
    foot = s0.shapes.add_textbox(Inches(1), Inches(6.5), Inches(4), Inches(0.5))
    foot.text_frame.paragraphs[0].add_run().text = "Globex Corp confidential"

    s1 = prs.slides.add_slide(blank)
    s1.shapes.add_picture(str(img), Inches(1), Inches(1), width=Inches(3))
    t = s1.shapes.add_table(2, 2, Inches(5), Inches(1), Inches(6), Inches(2)).table
    for ri in range(2):
        for ci in range(2):
            t.cell(ri, ci).text = "r%dc%d" % (ri, ci)

    s2 = prs.slides.add_slide(blank)
    tb2 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    tb2.text_frame.paragraphs[0].add_run().text = "Globex Corp roadmap"

    p = tmp_path / "deck.pptx"
    prs.save(p)
    return p


def shapes_of(data, idx):
    return {k: v for k, v in data["slides"][str(idx)].items() if not k.startswith("_")}


def find_sid(data, idx, **want):
    for sid, e in shapes_of(data, idx).items():
        if all(str(want[k]) in json.dumps(e.get(k, "")) for k in want):
            return sid
    raise AssertionError("no shape matching %s on slide %d" % (want, idx))


# ---------------------------------------------------------------------------


def test_docs_needs_no_file():
    r = run("docs")
    assert r.returncode == 0
    for section in ("GOLDEN RULES", "CREATE OPS", "TABLE STRUCTURE", "OUT OF SCOPE"):
        assert section in r.stdout


def test_inspect_reports_shapes_and_text(deck):
    data = inspect(deck, "--slide", "0")
    sid = find_sid(data, 0, paragraphs="Hello World")
    entry = shapes_of(data, 0)[sid]
    assert sid.startswith("s")
    assert entry["paragraphs"][0]["font_size"] == 30.0
    assert entry["paragraphs"][0]["bold"] is True


def test_set_text_inherits_formatting(deck, tmp_path):
    data = inspect(deck, "--slide", "0")
    sid = find_sid(data, 0, paragraphs="Hello World")
    out = tmp_path / "out.pptx"
    r = apply_patch(deck, [{"op": "set-text", "slide": 0, "shape": sid, "text": ["Replaced"]}], out)
    assert r.returncode == 0, r.stdout + r.stderr
    para = shapes_of(inspect(out, "--slide", "0"), 0)[sid]["paragraphs"][0]
    assert para["text"] == "Replaced"
    assert para["font_size"] == 30.0 and para["bold"] is True  # inherited


def test_replace_text_deck_wide_and_zero_hit_errors(deck, tmp_path):
    out = tmp_path / "out.pptx"
    r = apply_patch(deck, [{"op": "replace-text", "scope": "deck", "from": "Globex Corp", "to": "Acme"}], out)
    assert r.returncode == 0 and "2 occurrence(s)" in r.stdout
    r = apply_patch(deck, [{"op": "replace-text", "scope": "deck", "from": "Nonexistent", "to": "x"}], tmp_path / "o2.pptx")
    assert r.returncode != 0 and "not found" in r.stdout


def test_rejection_is_atomic_and_reports_all_errors(deck, tmp_path):
    before = deck.read_bytes()
    out = tmp_path / "never.pptx"
    r = apply_patch(deck, [
        {"op": "set-txt", "slide": 0, "shape": "s1", "text": ["x"]},
        {"op": "add-shape", "slide": 0, "kind": "blob", "at": [1, 1], "size": [2, 1]},
        {"op": "set-text", "slide": 0, "shape": "s9999", "text": ["x"]},
        {"op": "add-row", "slide": 99, "shape": "s1", "cells": ["a"]},
    ], out)
    assert r.returncode == 1
    assert "4 validation error(s)" in r.stdout
    assert "unknown op" in r.stdout
    assert "unknown shape kind" in r.stdout
    assert "shapes on slide 0" in r.stdout  # self-correcting: real inventory included
    assert "0-BASED" in r.stdout
    assert not out.exists()
    assert deck.read_bytes() == before


def test_runtime_failure_saves_nothing(deck, tmp_path):
    data = inspect(deck, "--slide", "1")
    table_sid = find_sid(data, 1, type="TABLE")
    out = tmp_path / "never.pptx"
    r = apply_patch(deck, [
        {"op": "set-text", "slide": 2, "shape": find_sid(data := inspect(deck, "--slide", "2"), 2, paragraphs="roadmap"), "text": ["changed"]},
        {"op": "add-row", "slide": 1, "shape": table_sid, "cells": ["only-one"]},
    ], out)
    assert r.returncode == 1 and "Nothing was saved" in r.stdout
    assert not out.exists()


def test_create_ops_and_style_by_name(deck, tmp_path, img):
    out = tmp_path / "out.pptx"
    r = apply_patch(deck, [
        {"op": "add-slide"},
        {"op": "add-shape", "slide": 3, "kind": "rect", "at": [0, 0], "size": [13.33, 7.5],
         "fill": "0F2B2E", "line": "none", "name": "bg"},
        {"op": "add-shape", "slide": 3, "kind": "textbox", "at": [1, 1], "size": [6, 1],
         "text": [{"text": "Built by an agent", "font_size": 32, "bold": True}], "name": "title"},
        {"op": "add-shape", "slide": 3, "kind": "CHEVRON", "at": [8, 1], "size": [2, 0.8], "fill": "4999A0"},
        {"op": "add-shape", "slide": 3, "kind": "line", "from": [1, 2.2], "to": [12, 2.2], "line_color": "C3DCDE"},
        {"op": "add-picture", "slide": 3, "image": str(img), "at": [9, 3], "width": 3},
        {"op": "add-table", "slide": 3, "at": [1, 3], "size": [7, 2], "rows": [["A", "B"], ["1", "2"]]},
        {"op": "set-style", "slide": 3, "shape": "title", "rotation": 6},  # created earlier in THIS patch
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    data = inspect(out, "--slide", "3")
    entries = shapes_of(data, 3)
    assert len(entries) == 6  # bg, title, chevron, line, picture, table
    title = entries[find_sid(data, 3, paragraphs="Built by an agent")]
    assert title["rotation"] == 6.0
    pic = entries[find_sid(data, 3, type="PICTURE")]
    assert pic["size"][0] == 3.0 and abs(pic["size"][1] - 1.8) < 0.05  # aspect preserved
    assert entries[find_sid(data, 3, type="TABLE")]["rows"] == [["A", "B"], ["1", "2"]]


def test_table_structure_ops(deck, tmp_path):
    sid = find_sid(inspect(deck, "--slide", "1"), 1, type="TABLE")
    out = tmp_path / "out.pptx"
    r = apply_patch(deck, [
        {"op": "add-row", "slide": 1, "shape": sid, "cells": ["r2c0", "r2c1"]},
        {"op": "add-col", "slide": 1, "shape": sid, "cells": ["c2r0", "c2r1", "c2r2"]},
        {"op": "delete-row", "slide": 1, "shape": sid, "row": 1},
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    rows = shapes_of(inspect(out, "--slide", "1"), 1)[sid]["rows"]
    assert rows == [["r0c0", "r0c1", "c2r0"], ["r2c0", "r2c1", "c2r2"]]


def test_merged_cell_guard(deck, tmp_path):
    from pptx.oxml.ns import qn
    prs = Presentation(deck)
    tbl = next(sh for sh in prs.slides[1].shapes if sh.has_table).table
    tbl._tbl.findall(qn("a:tr"))[0].findall(qn("a:tc"))[0].set("gridSpan", "2")
    merged = tmp_path / "merged.pptx"
    prs.save(merged)
    sid = find_sid(inspect(merged, "--slide", "1"), 1, type="TABLE")
    r = apply_patch(merged, [{"op": "add-row", "slide": 1, "shape": sid, "cells": ["a", "b"]}], tmp_path / "o.pptx")
    assert r.returncode == 1 and "merged cells" in r.stdout


def test_reorder_z(deck, tmp_path):
    data = inspect(deck, "--slide", "0")
    order = list(shapes_of(data, 0))
    out = tmp_path / "out.pptx"
    r = apply_patch(deck, [{"op": "reorder", "slide": 0, "shape": order[-1], "z": "back"}], out)
    assert r.returncode == 0
    assert list(shapes_of(inspect(out, "--slide", "0"), 0))[0] == order[-1]


def test_swap_image_media_global(deck, tmp_path):
    data = inspect(deck, "--slide", "1")
    media = shapes_of(data, 1)[find_sid(data, 1, type="PICTURE")]["media"]
    new_img = tmp_path / "new.png"
    Image.new("RGB", (200, 120), (200, 30, 30)).save(new_img)
    out = tmp_path / "out.pptx"
    r = apply_patch(deck, [{"op": "swap-image", "media": media, "image": str(new_img)}], out)
    assert r.returncode == 0 and "every reference changes" in r.stdout
    arc = "ppt/media/" + media
    assert zipfile.ZipFile(out).read(arc) != zipfile.ZipFile(deck).read(arc)
    r = apply_patch(deck, [{"op": "swap-image", "media": "nope.png", "image": str(new_img)}], tmp_path / "o2.pptx")
    assert r.returncode == 1 and "Available media" in r.stdout


def test_slides_subcommand_duplicates(deck, tmp_path):
    out = tmp_path / "out.pptx"
    r = run(deck, "slides", "0,2,2", "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    assert len(Presentation(out).slides._sldIdLst) == 3


def test_merge_subcommand(deck, tmp_path):
    out = tmp_path / "out.pptx"
    r = run(deck, "merge", deck, "--slides", "0", "--at", "1", "-o", out)
    assert r.returncode == 0, r.stdout + r.stderr
    assert len(Presentation(out).slides._sldIdLst) == 4


def test_fix_repairs_overflowing_text(deck, tmp_path):
    data = inspect(deck, "--slide", "0")
    sid = find_sid(data, 0, paragraphs="Hello World")
    grown = tmp_path / "grown.pptx"
    long_text = "This is a very long line of replacement text. " * 12
    r = apply_patch(deck, [{"op": "set-text", "slide": 0, "shape": sid, "text": [long_text]}], grown)
    assert r.returncode == 0
    r = run(grown, "fix", "--slides", "0", "--in-place")
    assert r.returncode == 0
    assert "fixed" in r.stdout


def test_diff_reports_changes(deck, tmp_path):
    out = tmp_path / "out.pptx"
    data = inspect(deck, "--slide", "0")
    sid = find_sid(data, 0, paragraphs="Hello World")
    apply_patch(deck, [
        {"op": "set-text", "slide": 0, "shape": sid, "text": ["Changed"]},
        {"op": "move", "slide": 0, "shape": sid, "by": [0.5, 0]},
    ], out)
    r = run(deck, "diff", out)
    assert "Hello World -> Changed" in r.stdout and "moved" in r.stdout


def test_replace_color_deck_wide(deck, tmp_path):
    out1 = tmp_path / "colored.pptx"
    r = apply_patch(deck, [
        {"op": "set-style", "slide": 0, "shape": find_sid(inspect(deck, "--slide", "0"), 0, paragraphs="Hello World"),
         "fill": "E8A33D", "color": "112233"},
    ], out1)
    assert r.returncode == 0, r.stdout + r.stderr
    out2 = tmp_path / "rethemed.pptx"
    r = apply_patch(out1, [
        {"op": "replace-color", "from": "E8A33D", "to": "0F5258"},
        {"op": "replace-color", "from": "#112233", "to": "FFFEFB"},  # leading # tolerated
    ], out2)
    assert r.returncode == 0, r.stdout + r.stderr
    assert "replace-color" in r.stdout and "E8A33D -> 0F5258" in r.stdout
    data = inspect(out2, "--slide", "0")
    blob = json.dumps(data)
    assert "0F5258" in blob and "FFFEFB" in blob
    assert "E8A33D" not in blob and "112233" not in blob


def test_replace_color_zero_hits_lists_palette(deck, tmp_path):
    out1 = tmp_path / "colored.pptx"
    sid = find_sid(inspect(deck, "--slide", "0"), 0, paragraphs="Hello World")
    apply_patch(deck, [{"op": "set-style", "slide": 0, "shape": sid, "fill": "AA00BB"}], out1)
    r = apply_patch(out1, [{"op": "replace-color", "from": "123456", "to": "0F5258"}], tmp_path / "no.pptx")
    assert r.returncode != 0
    assert "not found" in r.stdout and "AA00BB" in r.stdout  # error names the real palette
    # bad hex is caught at validation
    r = apply_patch(out1, [{"op": "replace-color", "from": "red", "to": "0F5258"}], tmp_path / "no2.pptx")
    assert r.returncode != 0 and "6-digit hex" in r.stdout


def test_lint_flags_text_under_picture(deck, tmp_path, img):
    # text first, picture later in z-order and intersecting → text renders
    # clipped behind the picture; apply must report covered_by
    out = tmp_path / "covered.pptx"
    r = apply_patch(deck, [
        {"op": "add-shape", "slide": 2, "kind": "textbox", "at": [1, 3], "size": [4, 1],
         "text": ["headline that gets trapped"], "name": "trapped"},
        {"op": "add-picture", "slide": 2, "image": str(img), "at": [2, 3.2], "width": 3},
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    assert "CLIPPED" in r.stdout
    data = inspect(out, "--slide", "2")
    entry = json.dumps(shapes_of(data, 2))
    assert "covered_by" in entry
    # fix reports it as residue with a z-order suggestion (never auto-resolved)
    r = run(out, "fix", "--slides", "2", "--in-place")
    assert "under PICTURE" in r.stdout and "z:front" in r.stdout


def test_text_on_top_of_picture_not_flagged(deck, tmp_path, img):
    # picture first, text later in z-order → text draws ON TOP: normal design
    out = tmp_path / "ontop.pptx"
    r = apply_patch(deck, [
        {"op": "add-picture", "slide": 2, "image": str(img), "at": [2, 3.2], "width": 3},
        {"op": "add-shape", "slide": 2, "kind": "textbox", "at": [1, 3], "size": [4, 1],
         "text": ["caption over the image"], "name": "caption"},
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    assert "covered_by" not in r.stdout and "CLIPPED" not in r.stdout


def test_overflowing_text_reaching_under_picture_is_flagged(deck, tmp_path, img):
    # the box itself clears the picture, but the ESTIMATED overflow runs under
    # it — the exact defect a thumbnail hides (re-wrapped last line clipped)
    out = tmp_path / "reach.pptx"
    long_text = "A very long serif headline that wraps far past its box. " * 10
    r = apply_patch(deck, [
        {"op": "add-shape", "slide": 2, "kind": "textbox", "at": [1, 1], "size": [3, 0.4],
         "text": [{"text": long_text, "font_size": 28}], "name": "spill"},
        {"op": "add-picture", "slide": 2, "image": str(img), "at": [1, 2.5], "width": 3},
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    assert "CLIPPED" in r.stdout


@pytest.mark.skipif(shutil.which("soffice") is None, reason="LibreOffice not installed")
def test_render_names_by_zero_based_index(deck, tmp_path):
    imgdir = tmp_path / "img"
    r = run(deck, "render", "-o", imgdir, "--slide", "2")
    assert r.returncode == 0, r.stdout + r.stderr
    assert (imgdir / "slide-2.jpg").exists()


def test_set_props_roundtrip_and_validation(deck, tmp_path):
    out = tmp_path / "props.pptx"
    r = apply_patch(deck, [
        {"op": "set-props", "title": "Q3 Review", "author": "Acme Strategy", "keywords": "ai, decks"},
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    data = inspect(out)
    assert data["props"]["title"] == "Q3 Review"
    assert data["props"]["author"] == "Acme Strategy"
    # unknown keys and empty op are validation errors
    r = apply_patch(deck, [{"op": "set-props", "ttile": "typo"}], tmp_path / "no.pptx")
    assert r.returncode != 0 and "unknown key" in r.stdout


def test_set_slide_hidden_and_solid_background(deck, tmp_path):
    out = tmp_path / "slideprops.pptx"
    r = apply_patch(deck, [
        {"op": "set-slide", "slide": 2, "hidden": True},
        {"op": "set-slide", "slide": 0, "background": "0F5258"},
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    data = inspect(out)
    assert data["hidden_slides"] == [2]
    xml_file = tmp_path / "s0.xml"
    run(out, "xml", "get", "--slide", "0", "-o", xml_file)
    xml = xml_file.read_text()
    assert "<p:bg>" in xml and "0F5258" in xml
    # unhide restores
    out2 = tmp_path / "unhidden.pptx"
    r = apply_patch(out, [{"op": "set-slide", "slide": 2, "hidden": False}], out2)
    assert r.returncode == 0
    assert "hidden_slides" not in inspect(out2)


def test_set_slide_image_background(deck, tmp_path, img):
    out = tmp_path / "bgimg.pptx"
    r = apply_patch(deck, [{"op": "set-slide", "slide": 1, "background": {"image": str(img)}}], out)
    assert r.returncode == 0, r.stdout + r.stderr
    xml_file = tmp_path / "s1.xml"
    run(out, "xml", "get", "--slide", "1", "-o", xml_file)
    xml = xml_file.read_text()
    assert "blipFill" in xml.split("spTree")[0]  # bg blipFill sits before the shape tree
    # file still opens and renders structure intact
    assert inspect(out)["slide_count"] == 3


def test_set_slide_transition_roundtrip(deck, tmp_path):
    out = tmp_path / "trans.pptx"
    r = apply_patch(deck, [
        {"op": "set-slide", "slide": 0, "transition": {"type": "fade", "speed": "slow", "advance_after": 5}},
        {"op": "set-slide", "slide": 1, "transition": {"type": "split", "orient": "vert", "dir": "out"}},
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    data = inspect(out)
    assert data["slides"]["0"]["_transition"] == {"type": "fade", "speed": "slow", "advance_after": 5.0}
    assert data["slides"]["1"]["_transition"] == {"type": "split", "orient": "vert", "dir": "out"}
    xml_file = tmp_path / "s0.xml"
    run(out, "xml", "get", "--slide", "0", "-o", xml_file)
    xml = xml_file.read_text()
    assert "<p:fade/>" in xml and 'spd="slow"' in xml and 'advTm="5000"' in xml
    # schema order: p:transition must come after cSld, never inside it
    assert xml.index("</p:cSld>") < xml.index("<p:transition")
    # transitions show up in diff (motion never shows in renders)
    d = run(deck, "diff", out)
    assert "transition" in d.stdout
    # setting again replaces, never stacks
    out2 = tmp_path / "trans2.pptx"
    r = apply_patch(out, [{"op": "set-slide", "slide": 0, "transition": {"type": "push", "dir": "u"}}], out2)
    assert r.returncode == 0
    assert inspect(out2)["slides"]["0"]["_transition"] == {"type": "push", "dir": "u"}
    run(out2, "xml", "get", "--slide", "0", "-o", xml_file)
    assert xml_file.read_text().count("<p:transition") == 1
    # "none" removes
    out3 = tmp_path / "trans3.pptx"
    r = apply_patch(out2, [{"op": "set-slide", "slide": 0, "transition": "none"}], out3)
    assert r.returncode == 0
    assert "_transition" not in inspect(out3)["slides"]["0"]


def test_set_slide_transition_validation(deck, tmp_path):
    r = apply_patch(deck, [
        {"op": "set-slide", "slide": 0, "transition": {"type": "swoosh"}},
        {"op": "set-slide", "slide": 1, "transition": {"type": "split", "dir": "l"}},
        {"op": "set-slide", "slide": 2, "transition": {"type": "fade", "speed": "ludicrous"}},
    ], tmp_path / "no.pptx")
    assert r.returncode != 0
    assert "swoosh" in r.stdout and "fade" in r.stdout  # unknown type teaches the vocabulary
    assert "invalid for split" in r.stdout and "in, out" in r.stdout
    assert "ludicrous" in r.stdout
    assert not (tmp_path / "no.pptx").exists()


ANIM_EFFECT_PAR = (
    '<p:par><p:cTn id="%(id1)d" fill="hold"><p:childTnLst>'
    '<p:set><p:cBhvr><p:cTn id="%(id2)d" dur="1" fill="hold"/>'
    '<p:tgtEl><p:spTgt spid="%(spid)s"/></p:tgtEl>'
    "<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>"
    '<p:to><p:strVal val="visible"/></p:to></p:set>'
    "</p:childTnLst></p:cTn></p:par>"
)
ANIM_TIMING = (
    '<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
    '<p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
    '<p:childTnLst><p:seq concurrent="1" nextAc="seek">'
    '<p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>%s</p:childTnLst></p:cTn>'
    "</p:seq></p:childTnLst></p:cTn></p:par></p:tnLst>"
    "<p:bldLst>%s</p:bldLst></p:timing>"
)


def _inject_animations(path, spids):
    """Give slide 0 a real PowerPoint timing tree: one entrance effect per shape id."""
    from pptx.oxml import parse_xml

    prs = Presentation(path)
    pars = "".join(
        ANIM_EFFECT_PAR % {"id1": 3 + 2 * i, "id2": 4 + 2 * i, "spid": spid}
        for i, spid in enumerate(spids)
    )
    blds = "".join('<p:bldP spid="%s" grpId="0"/>' % spid for spid in spids)
    prs.slides[0]._element.append(parse_xml(ANIM_TIMING % (pars, blds)))
    prs.save(path)


def test_delete_prunes_animation_refs(deck, tmp_path):
    data = inspect(deck)
    sid_a = find_sid(data, 0, paragraphs="Hello World")
    sid_b = find_sid(data, 0, paragraphs="confidential")
    _inject_animations(deck, [sid_a[1:], sid_b[1:]])
    out = tmp_path / "del1.pptx"
    r = apply_patch(deck, [{"op": "delete", "slide": 0, "shape": sid_a}], out)
    assert r.returncode == 0, r.stdout + r.stderr
    xml_file = tmp_path / "s0.xml"
    run(out, "xml", "get", "--slide", "0", "-o", xml_file)
    xml = xml_file.read_text()
    assert 'spid="%s"' % sid_a[1:] not in xml  # no dangling refs to the deleted shape
    assert 'spid="%s"' % sid_b[1:] in xml  # the other shape's animation survives
    # deleting the last animated shape leaves no timing tree at all
    out2 = tmp_path / "del2.pptx"
    r = apply_patch(out, [{"op": "delete", "slide": 0, "shape": sid_b}], out2)
    assert r.returncode == 0, r.stdout + r.stderr
    run(out2, "xml", "get", "--slide", "0", "-o", xml_file)
    assert "p:timing" not in xml_file.read_text()


def test_set_theme_colors_and_fonts(deck, tmp_path):
    import zipfile as zf

    out = tmp_path / "themed.pptx"
    r = apply_patch(deck, [
        {"op": "set-theme", "colors": {"accent1": "BB7B19", "dk1": "0F5258"},
         "fonts": {"major": "Georgia", "minor": "Verdana"}},
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    theme = zf.ZipFile(out).read("ppt/theme/theme1.xml").decode()
    assert 'val="BB7B19"' in theme and 'val="0F5258"' in theme
    assert 'typeface="Georgia"' in theme and 'typeface="Verdana"' in theme
    # bad slot name is a validation error that lists the real slots
    r = apply_patch(deck, [{"op": "set-theme", "colors": {"accent9": "112233"}}], tmp_path / "no.pptx")
    assert r.returncode != 0 and "accent6" in r.stdout


def test_alt_text_set_and_inspect(deck, tmp_path):
    data = inspect(deck, "--slide", "1")
    pic = find_sid(data, 1, type="PICTURE")
    out = tmp_path / "alt.pptx"
    r = apply_patch(deck, [
        {"op": "set-style", "slide": 1, "shape": pic, "alt_text": "bar chart of Q3 revenue"},
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    assert inspect(out, "--slide", "1")["slides"]["1"][pic]["alt_text"] == "bar chart of Q3 revenue"


def test_run_link_roundtrip(deck, tmp_path):
    data = inspect(deck, "--slide", "0")
    sid = find_sid(data, 0, paragraphs="Hello World")
    out = tmp_path / "linked.pptx"
    r = apply_patch(deck, [
        {"op": "set-text", "slide": 0, "shape": sid, "text": [
            {"runs": [{"text": "see "}, {"text": "the docs", "link": "https://example.com/docs"}]},
        ]},
    ], out)
    assert r.returncode == 0, r.stdout + r.stderr
    blob = json.dumps(inspect(out, "--slide", "0"))
    assert "https://example.com/docs" in blob


def test_notes_appear_in_inspect(deck, tmp_path):
    out = tmp_path / "noted.pptx"
    r = apply_patch(deck, [{"op": "set-notes", "slide": 0, "notes": "open with the Q3 number"}], out)
    assert r.returncode == 0
    assert inspect(out, "--slide", "0")["slides"]["0"]["_notes"] == "open with the Q3 number"


def test_set_notes_registers_notes_master(deck, tmp_path):
    # python-pptx creates the notesMaster part on first notes_slide access but
    # never adds p:notesMasterIdLst to presentation.xml. PowerPoint tolerates
    # the omission; Keynote rejects the entire file as "format is invalid".
    out = tmp_path / "noted.pptx"
    r = apply_patch(deck, [{"op": "set-notes", "slide": 0, "notes": "hello"}], out)
    assert r.returncode == 0
    with zipfile.ZipFile(out) as z:
        pres = z.read("ppt/presentation.xml").decode()
        rels = z.read("ppt/_rels/presentation.xml.rels").decode()
    assert "notesMasterIdLst" in pres
    m = re.search(r'<p:notesMasterId r:embed="(rId\d+)"|<p:notesMasterId r:id="(rId\d+)"', pres)
    assert m, "notesMasterId missing an r:id"
    rid = m.group(1) or m.group(2)
    assert rid in rels and "notesMaster" in rels
    # idempotent: a second set-notes must not add a duplicate list
    out2 = tmp_path / "noted2.pptx"
    r = apply_patch(out, [{"op": "set-notes", "slide": 1, "notes": "again"}], out2)
    assert r.returncode == 0
    with zipfile.ZipFile(out2) as z:
        assert z.read("ppt/presentation.xml").decode().count("notesMasterIdLst") == 2  # open+close tag
