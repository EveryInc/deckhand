"""End-to-end tests for html2patch: HTML -> patch -> deck.py apply -> inspect.

Requires playwright + chromium; every test is skipped when they're missing
(core hands-on-deck stays usable without them).
"""

import json
import shutil
import subprocess
import sys
from pathlib import Path

import pytest

SCRIPTS = Path(__file__).resolve().parent.parent / "skills" / "hands-on-deck" / "scripts"
H2P = SCRIPTS / "html2patch.py"
DECK = SCRIPTS / "deck.py"

playwright_ready = False
try:
    from playwright.sync_api import sync_playwright  # noqa: F401

    with sync_playwright() as _pw:
        _pw.chromium.launch().close()
    playwright_ready = True
except Exception:
    pass

pytestmark = pytest.mark.skipif(not playwright_ready, reason="playwright/chromium not available")

BODY = "body { width: 1280px; height: 720px; margin: 0; %s }"


def make_blank(path, slides=0):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = min(prs.slide_layouts, key=lambda l: len(l.shapes))
    for _ in range(slides):
        prs.slides.add_slide(blank)
    prs.save(path)
    return path


def h2p(*args):
    return subprocess.run(
        [sys.executable, str(H2P)] + [str(a) for a in args],
        capture_output=True, text=True,
    )


def deck(*args):
    return subprocess.run(
        [sys.executable, str(DECK)] + [str(a) for a in args],
        capture_output=True, text=True,
    )


def compile_html(tmp_path, html, *extra):
    page = tmp_path / "slide.html"
    page.write_text(html)
    deck_path = make_blank(tmp_path / "blank.pptx")
    out = tmp_path / "patch.json"
    r = h2p(page, "--deck", deck_path, "-o", out, *extra)
    assert r.returncode == 0, r.stderr
    return json.loads(out.read_text()), deck_path, out


def test_geometry_and_background(tmp_path):
    html = (
        "<html><head><style>" + BODY % "background:#112233;"
        + ".box { position:absolute; left:96px; top:192px; width:480px; height:96px;"
        "  background:#445566; }"
        "</style></head><body><div class='box'></div></body></html>"
    )
    patch, _, _ = compile_html(tmp_path, html)
    ops = patch["ops"]
    assert ops[0]["op"] == "add-slide"
    bg = ops[1]
    assert bg["fill"] == "112233" and bg["size"] == [13.333, 7.5]
    box = ops[2]
    assert box["kind"] == "rect" and box["fill"] == "445566"
    assert abs(box["at"][0] - 1.0) < 0.02 and abs(box["at"][1] - 2.0) < 0.02
    assert abs(box["size"][0] - 5.0) < 0.02 and abs(box["size"][1] - 1.0) < 0.02


def test_runs_transform_and_insets(tmp_path):
    html = (
        "<html><head><style>" + BODY % ""
        + "p { position:absolute; left:96px; top:96px; width:600px;"
        "  font: 24px Arial; color:#222222; padding: 12px 24px;"
        "  text-transform: uppercase; }"
        "b { color:#AA0000; }"
        "</style></head><body><p>mixed <b>weight</b> text</p></body></html>"
    )
    patch, _, _ = compile_html(tmp_path, html)
    text_ops = [o for o in patch["ops"] if o.get("kind") == "textbox"]
    assert len(text_ops) == 1
    op = text_ops[0]
    assert op["insets"] == [0.25, 0.125, 0.25, 0.125]  # padding 12px 24px
    runs = op["text"][0]["runs"]
    assert [r["text"] for r in runs] == ["MIXED ", "WEIGHT", " TEXT"]
    assert runs[1]["bold"] is True and runs[1]["color"] == "AA0000"
    assert "bold" not in runs[0] and runs[0]["color"] == "222222"
    assert runs[0]["font_size"] == 18.0  # 24px * 0.75


def test_nested_bullets_no_duplication(tmp_path):
    html = (
        "<html><head><style>" + BODY % ""
        + "ul, ol { position:absolute; left:96px; width:600px; font:16px Arial; }"
        "ul { top:96px; } ol { top:300px; }"
        "</style></head><body><ul><li>parent</li>"
        "<li>second<ul><li>child</li></ul></li></ul>"
        "<ol><li>step one</li><li>step two</li></ol></body></html>"
    )
    patch, _, _ = compile_html(tmp_path, html)
    boxes = [o for o in patch["ops"] if o.get("kind") == "textbox"]
    paras = boxes[0]["text"]
    texts = [(p.get("text") or p["runs"][0]["text"], p.get("level", 0)) for p in paras]
    assert texts == [("parent", 0), ("second", 0), ("child", 1)]
    assert all(p["bullet"] is True for p in paras)
    # ordered list compiles to auto-numbered paragraphs
    assert [p["bullet"] for p in boxes[1]["text"]] == ["number", "number"]


def test_widening_applies_to_all_blocks_and_serif_gets_more(tmp_path):
    # PPT redraws text wider than the browser; every block gets widened in its
    # anchor direction — 2% for sans faces, 4% for re-wrap-prone serif faces.
    html = (
        "<html><head><style>" + BODY % ""
        + "p { position:absolute; left:96px; width:600px; font:32px Georgia; }"
        ".serif { top:96px; }"
        ".sans  { top:400px; font-family: Arial; }"
        "</style></head><body>"
        "<p class='serif'>A long serif headline that wraps onto several lines here</p>"
        "<p class='sans'>A long sans headline that wraps onto several lines here</p>"
        "</body></html>"
    )
    patch, _, _ = compile_html(tmp_path, html)
    serif, sans = [o for o in patch["ops"] if o.get("kind") == "textbox"]
    # 600px = 6.25"; serif: ×1.04 = 6.5", sans: ×1.02 = 6.375"; left edge anchored
    assert abs(serif["size"][0] - 6.5) < 0.02, serif["size"]
    assert abs(sans["size"][0] - 6.375) < 0.02, sans["size"]
    assert abs(serif["at"][0] - 1.0) < 0.02 and abs(sans["at"][0] - 1.0) < 0.02


def test_object_fit_cover_and_figcaption(tmp_path):
    from PIL import Image

    Image.new("RGB", (300, 500), "#16777E").save(str(tmp_path / "tall.png"))
    html = (
        "<html><head><style>" + BODY % ""
        + "img { position:absolute; left:96px; top:96px; width:384px; height:192px;"
        "  object-fit: cover; }"
        "figcaption { position:absolute; left:96px; top:300px; font:14px Arial; color:#333333; }"
        "</style></head><body><img src='tall.png'>"
        "<figcaption>caption text here</figcaption></body></html>"
    )
    patch, _, _ = compile_html(tmp_path, html)
    pic = [o for o in patch["ops"] if o["op"] == "add-picture"][0]
    # 300x500 source into a 2:1 box: crop 35% top and bottom
    assert pic["crop"] == [0, 0.35, 0, 0.35]
    caps = [o for o in patch["ops"] if o.get("kind") == "textbox"]
    assert len(caps) == 1
    assert caps[0]["text"][0]["text"] == "caption text here"


def test_gradient_radius_and_partial_border(tmp_path):
    html = (
        "<html><head><style>" + BODY % ""
        + ".g { position:absolute; left:0; top:0; width:480px; height:240px;"
        "  background: linear-gradient(135deg, #102030, #405060);"
        "  border-radius: 24px; }"
        ".bar { position:absolute; left:600px; top:0; width:300px; height:100px;"
        "  border-left: 4px solid #AA0000; }"
        "</style></head><body><div class='g'></div><div class='bar'></div></body></html>"
    )
    patch, deck_path, patch_path = compile_html(tmp_path, html)
    grad = [o for o in patch["ops"] if "gradient" in o][0]
    assert grad["kind"] == "rounded_rect"
    assert grad["gradient"]["colors"] == ["102030", "405060"]
    # css 135deg (to bottom-right); pptx is counterclockwise from east
    assert grad["gradient"]["angle"] == 315
    assert grad["adjustments"] == [0.1]  # 24px / 240px
    assert grad["shadow"] is False  # theme shadows are suppressed
    lines = [o for o in patch["ops"] if o.get("kind") == "line"]
    assert len(lines) == 1
    assert lines[0]["line_color"] == "AA0000" and lines[0]["line_width"] == 3.0
    # the border-only bar's face stays hollow
    bar = [o for o in patch["ops"] if o.get("kind") == "rect" and o.get("fill") == "none"]
    assert len(bar) == 0 or all("gradient" not in o for o in bar)
    # and the whole thing applies cleanly
    r = deck(deck_path, "apply", patch_path, "-o", tmp_path / "out.pptx")
    assert r.returncode == 0, r.stdout + r.stderr


def test_table_neutralized_and_styled(tmp_path):
    html = (
        "<html><head><style>" + BODY % ""
        + "table { position:absolute; left:96px; top:96px; width:600px;"
        "  border-collapse:collapse; font:14px Arial; color:#333333; }"
        "th, td { height: 40px; } th { font-weight:bold; }"
        "th { background:#0B3D3A; color:#FFFFFF; }"
        "</style></head><body><table>"
        "<tr><th>A</th><th>B</th></tr><tr><td>1</td><td>2</td></tr>"
        "</table></body></html>"
    )
    patch, deck_path, patch_path = compile_html(tmp_path, html)
    tbl = [o for o in patch["ops"] if o["op"] == "add-table"][0]
    assert tbl["rows"] == [["A", "B"], ["1", "2"]]
    assert tbl["first_row"] is False and tbl["banding"] is False
    # mixed cell backgrounds become a per-cell fills grid
    assert tbl["fills"] == [["0B3D3A", "0B3D3A"], ["none", "none"]]
    assert len(tbl["col_widths"]) == 2 and abs(sum(tbl["col_widths"]) - 6.25) < 0.05
    cell_ops = {tuple(o["cell"]): o for o in patch["ops"] if o["op"] == "set-text"}
    assert cell_ops[(0, 0)]["text"][0]["bold"] and cell_ops[(0, 1)]["text"][0]["bold"]
    r = deck(deck_path, "apply", patch_path, "-o", tmp_path / "out.pptx")
    assert r.returncode == 0, r.stdout + r.stderr


def test_table_cell_br_keeps_words_apart(tmp_path):
    # textContent flattens <br> to nothing ("Who decideswhat runs next");
    # explicit breaks must survive as newlines
    html = (
        "<html><head><style>" + BODY % ""
        + "table { position:absolute; left:96px; top:96px; width:600px;"
        "  border-collapse:collapse; font:14px Arial; color:#333333; }"
        "td { height: 40px; }"
        "</style></head><body><table>"
        "<tr><td>Who decides<br>what runs next</td><td>The script</td></tr>"
        "</table></body></html>"
    )
    patch, deck_path, patch_path = compile_html(tmp_path, html)
    tbl = [o for o in patch["ops"] if o["op"] == "add-table"][0]
    assert tbl["rows"][0][0] == "Who decides\nwhat runs next"
    r = deck(deck_path, "apply", patch_path, "-o", tmp_path / "out.pptx")
    assert r.returncode == 0, r.stdout + r.stderr


def test_overflow_is_an_error(tmp_path):
    html = (
        "<html><head><style>" + BODY % ""
        + "div { height: 1000px; background:#222222; }"
        "</style></head><body><div></div></body></html>"
    )
    page = tmp_path / "over.html"
    page.write_text(html)
    r = h2p(page, "--size", "13.333x7.5", "--slide", "0")
    assert r.returncode == 1
    assert "overflows the body" in r.stderr


def test_compose_into_existing_deck(tmp_path):
    html = (
        "<html><head><style>" + BODY % ""
        + "p { position:absolute; left:96px; top:96px; font:24px Arial; }"
        "</style></head><body><p>new slide content</p></body></html>"
    )
    page = tmp_path / "slide.html"
    page.write_text(html)
    deck_path = make_blank(tmp_path / "deck.pptx", slides=2)
    patch_path = tmp_path / "patch.json"
    r = h2p(page, "--deck", deck_path, "-o", patch_path)
    assert r.returncode == 0, r.stderr
    patch = json.loads(patch_path.read_text())
    assert patch["ops"][0]["op"] == "add-slide"
    assert all(o.get("slide") == 2 for o in patch["ops"][1:])
    r = deck(deck_path, "apply", patch_path, "-o", tmp_path / "out.pptx")
    assert r.returncode == 0, r.stdout + r.stderr
    brief = deck(tmp_path / "out.pptx", "inspect", "--brief").stdout
    assert "slide 2" in brief and "new slide content" in brief


def test_zero_insets_no_false_overflow():
    """Regression: explicit zero text-frame margins must not be replaced by
    PowerPoint defaults in the overflow estimator (0 is falsy but real)."""
    sys.path.insert(0, str(SCRIPTS))
    try:
        import importlib
        import inventory
        importlib.reload(inventory)
    finally:
        sys.path.pop(0)

    class Frame:
        margin_left = margin_top = margin_right = margin_bottom = 0

    class Probe(inventory.ShapeData):
        def __init__(self):  # bypass full init; only geometry matters here
            self.width, self.height = 2.0, 0.5

    usable_w, usable_h = Probe()._get_usable_dimensions(Frame())
    assert usable_h == pytest.approx(0.5 * 96)  # full height, no phantom margins
    assert usable_w == pytest.approx(2.0 * 96)


def test_anchor_becomes_hyperlink_run(tmp_path):
    html = (
        "<html><head><style>" + BODY % ""
        + "p { position:absolute; left:96px; top:96px; width:600px; font:20px Arial; }"
        "a { color:#BB7B19; }"
        "</style></head><body>"
        "<p>read <a href='https://example.com/docs'>the docs</a> first</p>"
        "</body></html>"
    )
    patch, deck_path, patch_path = compile_html(tmp_path, html)
    op = [o for o in patch["ops"] if o.get("kind") == "textbox"][0]
    runs = op["text"][0]["runs"]
    linked = [r for r in runs if r.get("link")]
    assert len(linked) == 1
    assert linked[0]["text"] == "the docs" and linked[0]["link"] == "https://example.com/docs"
    r = deck(deck_path, "apply", patch_path, "-o", tmp_path / "out.pptx")
    assert r.returncode == 0, r.stdout + r.stderr
